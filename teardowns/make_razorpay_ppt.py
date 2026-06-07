"""
Razorpay Product Teardown — 20-Slide Deck
Shraddha Singh · June 2026

Brand: Razorpay
  Primary dark : #012652  (Prussian Blue)
  Accent blue  : #0D94FB  (Dodger Blue)
  Purple accent: #6822CC
  Font family  : Mulish → Helvetica fallback
"""

import os
from io import BytesIO
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Brand Colors ─────────────────────────────────────────────────────────
RPY_NAVY    = RGBColor(0x01, 0x26, 0x52)   # Prussian Blue — primary bg
RPY_DARK    = RGBColor(0x01, 0x0F, 0x28)   # Darker variant for slide bg
RPY_CARD    = RGBColor(0x03, 0x1D, 0x42)   # Card surfaces
RPY_BLUE    = RGBColor(0x0D, 0x94, 0xFB)   # Dodger Blue — accent
RPY_PURPLE  = RGBColor(0x68, 0x22, 0xCC)   # Deep Purple — secondary accent
RPY_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
RPY_GRAY    = RGBColor(0x8A, 0xA5, 0xBF)   # Muted text
RPY_LGRAY   = RGBColor(0xC8, 0xD8, 0xE8)   # Light muted text
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)
ACCENT_AMBER = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT_RED   = RGBColor(0xEF, 0x44, 0x44)

ASSETS = "/tmp/razorpay_assets"
os.makedirs(ASSETS, exist_ok=True)

FONT_PATH = "/System/Library/Fonts/Helvetica.ttc"


def _font(size):
    try:
        return ImageFont.truetype(FONT_PATH, size=size)
    except Exception:
        return ImageFont.load_default()


# ════════════════════════════════════════════════════════════════════════
# IMAGE ASSET GENERATORS
# ════════════════════════════════════════════════════════════════════════

def make_rpy_logo(filepath=None):
    """Razorpay wordmark — navy bg, blue accent bar, white text."""
    W, H = 480, 160
    img = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)
    draw.rounded_rectangle([0, 0, W-1, H-1], radius=16,
                            fill=(1, 15, 40, 255))
    # Blue accent bar top
    draw.rounded_rectangle([0, 0, W-1, 8], radius=4,
                            fill=(13, 148, 251, 255))
    # "razorpay" wordmark
    f = _font(60)
    text = "razorpay"
    bbox = draw.textbbox((0, 0), text, font=f)
    tw = bbox[2] - bbox[0]
    th = bbox[3] - bbox[1]
    # Blue "r" then white rest
    draw.text(((W - tw) // 2, (H - th) // 2 + 6), text[0],
              font=f, fill=(13, 148, 251, 255))
    offset_x = (W - tw) // 2 + draw.textbbox((0, 0), text[0], font=f)[2]
    draw.text((offset_x, (H - th) // 2 + 6), text[1:],
              font=f, fill=(255, 255, 255, 255))
    path = filepath or f"{ASSETS}/razorpay_logo.png"
    img.save(path)
    return path


def make_badge(text, bg_hex, fg_hex="FFFFFF", size=180, filepath=None):
    img = Image.new("RGBA", (size, size), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)
    bg = tuple(int(bg_hex[i:i+2], 16) for i in (0, 2, 4)) + (255,)
    fg = tuple(int(fg_hex[i:i+2], 16) for i in (0, 2, 4)) + (255,)
    draw.rounded_rectangle([0, 0, size-1, size-1], radius=26, fill=bg)
    f = _font(size // 3)
    bbox = draw.textbbox((0, 0), text, font=f)
    tw = bbox[2] - bbox[0]; th = bbox[3] - bbox[1]
    draw.text(((size - tw) // 2, (size - th) // 2 - 4), text, font=f, fill=fg)
    path = filepath or f"{ASSETS}/{text.lower().replace(' ', '_')}.png"
    img.save(path)
    return path


def make_revenue_chart(filepath=None):
    """Bar chart — Razorpay revenue FY21–FY24."""
    W, H = 720, 320
    img = Image.new("RGBA", (W, H), (3, 29, 66, 255))
    draw = ImageDraw.Draw(img)
    f_hd = _font(20); f_sm = _font(15)
    blue  = (13, 148, 251, 255)
    gray  = (138, 165, 191, 255)
    white = (255, 255, 255, 255)
    draw.text((20, 14), "Revenue (₹ Crore)  ·  FY21–FY24", font=f_hd,
              fill=(13, 148, 251, 255))
    data = [
        ("FY 2021",  390,  "345670"),
        ("FY 2022", 1346,  "235890"),
        ("FY 2023", 2285,  "134CAF"),
        ("FY 2024", 2600,  "0D94FB"),
    ]
    bar_x = 120; max_bar = W - bar_x - 110; bar_h = 44; gap = 18
    max_val = max(v for _, v, _ in data)
    y = 56
    for label, val, col in data:
        ratio = val / max_val
        bc = tuple(int(col[i:i+2], 16) for i in (0, 2, 4)) + (255,)
        draw.rectangle([bar_x, y, bar_x + int(max_bar*ratio), y+bar_h], fill=bc)
        draw.text((14, y + 12), label, font=f_sm, fill=white)
        draw.text((bar_x + int(max_bar*ratio) + 8, y + 12),
                  f"₹{val:,}", font=f_sm, fill=(13, 148, 251, 255))
        y += bar_h + gap
    path = filepath or f"{ASSETS}/rpy_revenue.png"
    img.save(path)
    return path


def make_stack_diagram(filepath=None):
    """Razorpay product stack — 4 layers stacked."""
    W, H = 680, 360
    img = Image.new("RGBA", (W, H), (3, 29, 66, 255))
    draw = ImageDraw.Draw(img)
    f_lg = _font(18); f_sm = _font(13)
    layers = [
        ("Lending  (Razorpay Capital)", "0D94FB"),
        ("Banking  (RazorpayX + Opfin)", "1880D0"),
        ("Collections  (Smart Collect · Links · Pages)", "0F60A0"),
        ("Payments  (Gateway · SDK · Checkout)", "0A3A6A"),
    ]
    block_h = 72; gap = 10; y = 20
    for label, col in layers:
        bc = tuple(int(col[i:i+2], 16) for i in (0, 2, 4)) + (255,)
        draw.rounded_rectangle([30, y, W-30, y+block_h], radius=10, fill=bc)
        draw.text((54, y + block_h//2 - 10), label, font=f_lg, fill=(255, 255, 255, 255))
        y += block_h + gap
    draw.text((30, H-28), "← Foundation layer       High-margin growth →",
              font=f_sm, fill=(138, 165, 191, 255))
    path = filepath or f"{ASSETS}/rpy_stack.png"
    img.save(path)
    return path


def make_upi_market(filepath=None):
    """UPI market share donut placeholder."""
    W, H = 600, 380
    img = Image.new("RGBA", (W, H), (3, 29, 66, 255))
    draw = ImageDraw.Draw(img)
    f_lg = _font(32); f_sm = _font(16); f_xs = _font(13)
    # Concentric rings representing market
    draw.ellipse([100, 30, 500, 350], outline=(13, 148, 251, 80), width=2)
    draw.ellipse([160, 80, 440, 300], outline=(13, 148, 251, 140), width=2)
    draw.ellipse([210, 120, 390, 260], fill=(13, 42, 85, 255),
                 outline=(13, 148, 251, 255), width=2)
    draw.text((240, 162), "8M+", font=f_lg, fill=(13, 148, 251, 255))
    draw.text((228, 205), "Merchants", font=f_sm, fill=(138, 165, 191, 255))
    draw.text((20, 45), "~64M\nIndia\nSMEs", font=f_xs, fill=(138, 165, 191, 255))
    draw.text((490, 45), "~8M\nRazorpay\nMerchants", font=f_xs, fill=(13, 148, 251, 255))
    draw.text((20, 310), "TAM: ₹12T+ payments", font=f_xs, fill=(138, 165, 191, 255))
    path = filepath or f"{ASSETS}/rpy_market.png"
    img.save(path)
    return path


def load_or_badge(real_path, badge_text, badge_bg, badge_fg="FFFFFF", size=200):
    """Use the real logo PNG if it exists; fall back to a styled badge."""
    if os.path.exists(real_path):
        try:
            img = Image.open(real_path).convert("RGBA")
            # Compose onto navy bg to match slide aesthetic
            bg = Image.new("RGBA", (size, size), (1, 15, 40, 255))
            scale = min((size-16)/img.width, (size-16)/img.height)
            nw, nh = int(img.width*scale), int(img.height*scale)
            img = img.resize((nw, nh), Image.LANCZOS)
            bg.paste(img, ((size-nw)//2, (size-nh)//2), img)
            out = real_path.replace("_real.png", "_final.png")
            bg.convert("RGB").save(out)
            return out
        except Exception:
            pass
    return make_badge(badge_text, badge_bg, badge_fg, size=size)


# ── Generate assets ──────────────────────────────────────────────────────
print("Generating image assets...")
RPY_LOGO      = make_rpy_logo()
REV_CHART     = make_revenue_chart()
STACK_DIAGRAM = make_stack_diagram()
MARKET_MAP    = make_upi_market()

# Real logos (GitHub org avatars) with badge fallbacks
LOGO_RAZORPAY = load_or_badge(f"{ASSETS}/razorpay_real.png", "rpy",     "012652", "0D94FB")
LOGO_PAYU     = load_or_badge(f"{ASSETS}/payu_real.png",     "PayU",    "0070CC")
LOGO_CASHFREE = load_or_badge(f"{ASSETS}/cashfree_real.png", "CF",      "00A86B")
LOGO_STRIPE   = load_or_badge(f"{ASSETS}/stripe_real.png",   "stripe",  "6772E5")
LOGO_PHONEPE  = load_or_badge(f"{ASSETS}/phonepe_real.png",  "PP",      "5F259F")
LOGO_JUSPAY   = load_or_badge(f"{ASSETS}/juspay_real.png",   "JP",      "1A5276")
LOGO_CCAVENUE = make_badge("CCA", "CC3300", filepath=f"{ASSETS}/ccavenue_final.png")
print("Assets ready.")


# ════════════════════════════════════════════════════════════════════════
# PPT HELPERS
# ════════════════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_pt=0):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.line.width = Pt(line_pt)
    if fill_rgb:
        s.fill.solid(); s.fill.fore_color.rgb = fill_rgb
    else:
        s.fill.background()
    if line_rgb:
        s.line.color.rgb = line_rgb
    else:
        s.line.fill.background()
    return s


def add_text(slide, text, l, t, w, h, size=14, bold=False,
             color=None, align=PP_ALIGN.LEFT, italic=False):
    if color is None:
        color = RPY_WHITE
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tb


def add_ml(slide, lines, l, t, w, h, size=12, default_color=None):
    if default_color is None:
        default_color = RPY_WHITE
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True; tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            txt, bld, clr = item, False, default_color
        else:
            txt = item[0]
            bld = item[1] if len(item) > 1 else False
            clr = item[2] if len(item) > 2 else default_color
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(3)
        r = p.add_run(); r.text = txt
        r.font.size = Pt(size); r.font.bold = bld; r.font.color.rgb = clr


def add_img(slide, path, l, t, w, h=None):
    try:
        if h:
            slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
        else:
            slide.shapes.add_picture(path, Inches(l), Inches(t), width=Inches(w))
    except Exception as e:
        print(f"  img skip: {e}")


def dark_bg(slide):
    """Navy background + thin blue accent bars top & bottom."""
    add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=RPY_DARK)
    add_rect(slide, 0, 0,    13.33, 0.05, fill_rgb=RPY_BLUE)
    add_rect(slide, 0, 7.45, 13.33, 0.05, fill_rgb=RPY_BLUE)


def section_label(slide, tag, subtitle=""):
    add_text(slide, tag,      0.6, 0.18, 12, 0.35, size=9.5,
             color=RPY_BLUE, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.6, 0.48, 12, 0.3, size=8.5, color=RPY_GRAY)


def new_slide(tag=None, subtitle=None):
    s = prs.slides.add_slide(BLANK)
    dark_bg(s)
    if tag:
        section_label(s, tag, subtitle)
    return s


def section_divider(label, subtitle=""):
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=RPY_NAVY)
    # Left blue bar
    add_rect(slide, 0, 0, 0.14, 7.5, fill_rgb=RPY_BLUE)
    add_rect(slide, 0, 0, 13.33, 0.05, fill_rgb=RPY_BLUE)
    add_rect(slide, 0, 7.45, 13.33, 0.05, fill_rgb=RPY_BLUE)
    add_text(slide, label, 0.55, 2.8, 12, 1.2, size=50, bold=True,
             color=RPY_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.55, 4.2, 12, 0.5, size=16,
                 color=RPY_BLUE, align=PP_ALIGN.LEFT)
    add_text(slide, "Razorpay  ·  Product Teardown  ·  Shraddha Singh  ·  June 2026",
             0.55, 7.0, 12, 0.38, size=9, color=RPY_GRAY)
    return slide


def footer_strip(slide, text, highlight=None):
    add_rect(slide, 0.4, 6.5, 12.5, 0.88, fill_rgb=RPY_CARD)
    if highlight:
        add_rect(slide, 0.4, 6.5, 0.06, 0.88, fill_rgb=RPY_BLUE)
    add_text(slide, text, 0.62, 6.58, 12.0, 0.72, size=11,
             italic=True, color=RPY_LGRAY)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=RGBColor(0x00, 0x0C, 0x20))
add_rect(slide, 0, 0,    13.33, 0.07, fill_rgb=RPY_BLUE)
add_rect(slide, 0, 7.43, 13.33, 0.07, fill_rgb=RPY_BLUE)
# Right panel
add_rect(slide, 7.3, 0, 6.03, 7.5, fill_rgb=RPY_NAVY)
# Diagonal gradient hint via purple rect
add_rect(slide, 8.5, 3.5, 4.5, 4.0, fill_rgb=RGBColor(0x10, 0x04, 0x24))
add_img(slide, LOGO_RAZORPAY if os.path.exists(str(LOGO_RAZORPAY)) else RPY_LOGO, 7.6, 1.8, 4.8, 1.6)
# Left: Title
add_text(slide, "razorpay", 0.9, 1.0, 7, 2.0, size=90, bold=True, color=RPY_WHITE)
add_text(slide, "Product Teardown", 0.9, 3.25, 7, 0.65, size=28, color=RPY_BLUE)
add_text(slide, "From payment gateway to India's financial OS —\nand the PM questions behind every product bet",
         0.9, 4.05, 6.8, 0.75, size=13.5, color=RPY_GRAY)
add_rect(slide, 0.9, 5.1, 2.8, 0.04, fill_rgb=RPY_BLUE)
add_text(slide, "Shraddha Singh", 0.9, 5.22, 5.5, 0.38, size=14, bold=True)
add_text(slide, "Product Manager  ·  Fintech & Enterprise SaaS",
         0.9, 5.65, 6.5, 0.35, size=12, color=RPY_GRAY)
add_text(slide, "June 2026", 0.9, 6.1, 3, 0.32, size=11, color=RPY_BLUE)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ════════════════════════════════════════════════════════════════════════
slide = new_slide("AGENDA", "What this teardown covers")
add_text(slide, "Agenda", 0.6, 0.9, 12, 0.65, size=34, bold=True)

sections = [
    ("01", "Context",          "What is Razorpay · The B2B user · Core problems solved"),
    ("02", "The Product",      "Platform evolution · Product suite · Feature analysis"),
    ("03", "The Business",     "3-layer business model · Revenue data · Market opportunity"),
    ("04", "The Competition",  "Razorpay vs Stripe · PayU · Cashfree · PhonePe for Business"),
    ("05", "Product Thinking", "3 product opportunities · PM Takeaways"),
]
for i, (num, title, detail) in enumerate(sections):
    if i < 3:
        x, y, w = 0.5 + i * 4.28, 1.85, 4.0
    else:
        x, y, w = 0.5 + (i-3) * 6.45, 4.55, 6.0
    add_rect(slide, x, y, w, 2.2, fill_rgb=RPY_CARD)
    add_rect(slide, x, y, w, 0.05, fill_rgb=RPY_BLUE)
    add_text(slide, num,    x+0.2, y+0.12, 0.8, 0.5, size=22, bold=True, color=RPY_BLUE)
    add_text(slide, title,  x+0.2, y+0.65, w-0.4, 0.5, size=14, bold=True)
    add_text(slide, detail, x+0.2, y+1.18, w-0.4, 0.9, size=11, color=RPY_GRAY)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 3 — WHAT IS RAZORPAY
# ════════════════════════════════════════════════════════════════════════
slide = new_slide("WHAT IS RAZORPAY", "Overview")
add_text(slide, "What is Razorpay?", 0.6, 0.9, 12, 0.65, size=34, bold=True)
add_rect(slide, 0.6, 1.72, 0.05, 0.95, fill_rgb=RPY_BLUE)
add_text(slide, "\"They won because their API documentation was better —\n"
         "and in a developer-bought product, documentation is the product.\"",
         0.85, 1.72, 8.5, 0.55, size=14, italic=True, color=RPY_BLUE)
bullets = [
    "Founded 2014 by Harshil Mathur & Shashvat Nakrani — IIT Roorkee → YCombinator",
    "Started as a payment gateway API; now: payments, banking, payroll, and lending",
    "Won the market not on price but on developer experience — best API docs in the industry",
    "Platform expansion: each new product uses data from the previous one",
    "Positioning: India's financial operating system for SMEs and enterprises",
]
add_ml(slide, [("• " + b, False, RPY_WHITE) for b in bullets],
       0.85, 2.4, 7.8, 2.8, size=13.5)
add_img(slide, LOGO_RAZORPAY if os.path.exists(str(LOGO_RAZORPAY)) else RPY_LOGO, 9.6, 1.85, 3.3, 1.1)

kpi = [("8M+", "Businesses on platform"),
       ("$90B+", "Annual payments processed"),
       ("$7.5B", "Valuation (Series F 2021)"),
       ("100+", "Payment methods")]
for i, (num, label) in enumerate(kpi):
    kx = 0.6 + i * 3.12
    add_rect(slide, kx, 5.38, 2.98, 0.92, fill_rgb=RPY_CARD)
    add_rect(slide, kx, 5.38, 2.98, 0.05, fill_rgb=RPY_BLUE)
    add_text(slide, num,   kx+0.18, 5.48, 2.7, 0.44, size=20, bold=True, color=RPY_BLUE)
    add_text(slide, label, kx+0.18, 5.94, 2.7, 0.3,  size=10, color=RPY_GRAY)

footer_strip(slide, "The simplest version: won the gateway market on developer experience, "
             "then used that distribution to expand into every financial product an SME needs.",
             highlight=True)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 4 — THE B2B USER
# ════════════════════════════════════════════════════════════════════════
slide = new_slide("THE B2B USER", "Three distinct user types — very different needs")
add_text(slide, "The B2B User — Why Complexity Matters",
         0.6, 0.9, 12, 0.65, size=30, bold=True)

user_data = [
    ("Developer\n/ Tech Lead",
     "Builds checkout & payment flows; evaluates APIs",
     "Clean docs, reliable uptime, instant sandbox, fast integration. "
     "Will internally advocate for the tool that makes them look good.",
     "Initiates adoption",
     RPY_BLUE),
    ("Finance\n/ Operations",
     "Owns reconciliation, settlements, compliance",
     "Accurate reporting, predictable settlement times, audit trail, "
     "GST handling. Will block adoption if reporting doesn't work.",
     "Veto power",
     RPY_PURPLE),
    ("Business Owner\n/ Founder",
     "Cares about conversion, revenue, growth",
     "Higher checkout success rate, lower payment failure rate, access "
     "to credit when needed. Buys on outcomes, not features.",
     "Writes the cheque",
     ACCENT_GREEN),
]
for i, (utype, role, want, power, col) in enumerate(user_data):
    x = 0.45 + i * 4.25
    add_rect(slide, x, 1.85, 4.05, 4.45, fill_rgb=RPY_CARD)
    add_rect(slide, x, 1.85, 4.05, 0.05, fill_rgb=col)
    add_text(slide, utype, x+0.2, 1.94, 3.7, 0.72, size=14, bold=True, color=col)
    add_text(slide, role,  x+0.2, 2.7,  3.7, 0.48, size=11, italic=True, color=RPY_LGRAY)
    add_rect(slide, x+0.2, 3.25, 3.6, 0.03, fill_rgb=RPY_CARD)
    add_text(slide, want,  x+0.2, 3.35, 3.6, 2.18, size=11.5, color=RPY_WHITE)
    # Power badge at bottom
    add_rect(slide, x+0.2, 5.88, 3.6, 0.32, fill_rgb=RGBColor(0x02, 0x20, 0x44))
    add_text(slide, power, x+0.3, 5.91, 3.4, 0.26, size=10, bold=True, color=col)

footer_strip(slide, "PM insight: Design first for the developer (initiates) while ensuring the "
             "finance team can't block. The developer self-serves; the finance team has veto power. You need both.",
             highlight=True)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 5 — CORE PROBLEMS SOLVED
# ════════════════════════════════════════════════════════════════════════
slide = new_slide("CORE PROBLEMS SOLVED", "What India's payment landscape looked like before Razorpay")
add_text(slide, "The Problems Razorpay Solved", 0.6, 0.9, 12, 0.65, size=32, bold=True)
problems = [
    ("Gateway integration took weeks",
     "Banks and CCAvenue required manual forms, verification, enterprise sales cycles. "
     "A startup couldn't go live in a day.",
     "→ Self-serve in 15 min"),
    ("API documentation was unusable",
     "Competitors had PDFs, SOAP APIs, guides written by lawyers. "
     "Developers copied Stack Overflow to understand basic auth flows.",
     "→ Developer-first docs"),
    ("Payment failure rates: 65–70%",
     "Industry average success rate was ~65–70%. Lost transactions = lost revenue. "
     "Nobody treated this as a PM metric.",
     "→ 85%+ success rate"),
    ("SMEs had no real banking product",
     "Traditional SME accounts were enterprise accounts stripped down. "
     "No API access, no automation, no real-time visibility.",
     "→ RazorpayX banking API"),
]
for i, (prob, insight, fix) in enumerate(problems):
    x = 0.45 + (i % 2) * 6.45
    y = 1.85 + (i // 2) * 2.4
    add_rect(slide, x, y, 6.15, 2.15, fill_rgb=RPY_CARD)
    add_rect(slide, x, y, 6.15, 0.05, fill_rgb=RPY_BLUE)
    add_text(slide, f"0{i+1}", x+0.2, y+0.1, 0.7, 0.55, size=22, bold=True, color=RPY_BLUE)
    add_text(slide, prob,    x+0.2, y+0.65, 5.7, 0.45, size=13, bold=True)
    add_text(slide, insight, x+0.2, y+1.12, 5.7, 0.85, size=11, color=RPY_GRAY)
    add_rect(slide, x+0.2, y+1.88, 5.6, 0.22, fill_rgb=RGBColor(0x02, 0x20, 0x44))
    add_text(slide, fix,     x+0.3, y+1.9,  5.4, 0.18, size=10, bold=True, color=RPY_BLUE)


# ════════════════════════════════════════════════════════════════════════
# SECTION DIVIDER — THE PRODUCT
# ════════════════════════════════════════════════════════════════════════
section_divider("The Product", "Platform Evolution · Suite · Feature Analysis")


# ════════════════════════════════════════════════════════════════════════
# SLIDE 7 — PRODUCT EVOLUTION TIMELINE
# ════════════════════════════════════════════════════════════════════════
slide = new_slide("PRODUCT EVOLUTION", "From API to full financial OS — 10 years of platform expansion")
add_text(slide, "How the Platform Evolved", 0.6, 0.9, 12, 0.65, size=32, bold=True)

timeline = [
    ("2014", "Gateway",    "Payments API\n+ sandbox"),
    ("2016–17", "Distribution", "PCI-DSS\nint'l cards"),
    ("2018–19", "Product Depth", "Links · Pages\nSubscriptions"),
    ("2020", "Adjacency", "RazorpayX\nSmart Collect"),
    ("2021", "Full Stack", "Capital\n+ Opfin"),
    ("2022–24", "Enterprise", "Int'l payments\nAI fraud layer"),
]
dot_y = 2.6
add_rect(slide, 0.55, dot_y+0.18, 12.2, 0.04, fill_rgb=RPY_BLUE)
xs = [0.55, 2.6, 4.65, 6.7, 8.75, 10.8]

for i, (yr, lbl, desc) in enumerate(timeline):
    x = xs[i]
    dot = slide.shapes.add_shape(9, Inches(x+0.78), Inches(dot_y+0.1),
                                  Inches(0.28), Inches(0.28))
    dot.fill.solid(); dot.fill.fore_color.rgb = RPY_BLUE
    dot.line.fill.background()
    add_text(slide, yr,   x, dot_y-0.98, 2.0, 0.38, size=11, bold=True,
             color=RPY_BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, lbl,  x, dot_y-0.55, 2.0, 0.38, size=10,
             color=RPY_GRAY, align=PP_ALIGN.CENTER)
    add_rect(slide, x+0.15, dot_y+0.52, 1.6, 1.5, fill_rgb=RPY_CARD)
    add_text(slide, desc, x+0.2, dot_y+0.58, 1.55, 1.4, size=10,
             align=PP_ALIGN.CENTER, color=RPY_WHITE)

add_rect(slide, 0.4, 4.95, 12.5, 0.88, fill_rgb=RPY_CARD)
add_rect(slide, 0.4, 4.95, 0.06, 0.88, fill_rgb=RPY_BLUE)
add_text(slide, "THE SEQUENCING LOGIC:",
         0.6, 5.03, 2.8, 0.3, size=9, bold=True, color=RPY_BLUE)
add_text(slide, "Payment collection → Banking → Payroll → Lending.  Each step used data from the previous one.  "
         "Banking gave cash flow visibility.  Payroll gave employee/burn data.  Both made Razorpay Capital's "
         "underwriting structurally better than any traditional lender could achieve.",
         0.6, 5.36, 12.1, 0.42, size=11, color=RPY_WHITE)

milestones = [("2014", "YC batch — gateway launch"),
              ("2020", "RazorpayX — banking API"),
              ("2021", "Capital launched — lending enters")]
for i, (yr, lbl) in enumerate(milestones):
    x = 0.4 + i * 4.3
    add_rect(slide, x, 6.05, 4.1, 0.65, fill_rgb=RGBColor(0x01, 0x1A, 0x38))
    add_rect(slide, x, 6.05, 4.1, 0.04, fill_rgb=RPY_BLUE)
    add_text(slide, yr,  x+0.18, 6.12, 1.0, 0.32, size=12, bold=True, color=RPY_BLUE)
    add_text(slide, lbl, x+1.28, 6.12, 2.7, 0.44, size=11)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 8 — PRODUCT STACK VISUAL
# ════════════════════════════════════════════════════════════════════════
slide = new_slide("PRODUCT SUITE", "The 4-layer platform stack")
add_text(slide, "The Razorpay Product Stack", 0.6, 0.9, 12, 0.65, size=32, bold=True)
add_img(slide, STACK_DIAGRAM, 0.4, 1.6, 6.0, 3.4)

products = [
    ("Payment Gateway",    "MDR per transaction",      "Foundation; self-serve in 15 min; 85%+ success rate"),
    ("RazorpayX",         "Float income + SaaS",       "Current accounts with API; corporate cards; vendor payments"),
    ("Opfin (Payroll)",   "Per-employee SaaS",         "Payroll, PF, ESI, TDS — acquired ₹125 Cr (2021)"),
    ("Razorpay Capital",  "Interest + processing fee", "SME loans underwritten via transaction history — no bureau"),
    ("Smart Collect",     "Per-transaction fee",       "UPI QR, virtual accounts, B2B collections"),
    ("Links & Pages",     "MDR per transaction",       "No-code payment collection for non-technical merchants"),
]
add_rect(slide, 6.7, 1.62, 6.2, 0.38, fill_rgb=RPY_BLUE)
add_text(slide, "  Product · Revenue Model · PM Insight",
         6.7, 1.62, 6.2, 0.38, size=11, bold=True, color=RPY_DARK)
for i, (prod, rev, insight) in enumerate(products):
    ry = 2.04 + i * 0.6
    bg = RPY_CARD if i % 2 == 0 else RGBColor(0x01, 0x1A, 0x38)
    add_rect(slide, 6.7, ry, 6.2, 0.58, fill_rgb=bg)
    add_text(slide, prod,    6.88, ry+0.06, 1.7,  0.42, size=11, bold=True, color=RPY_BLUE)
    add_text(slide, rev,     8.62, ry+0.06, 1.85, 0.42, size=10, color=RPY_LGRAY)
    add_text(slide, insight, 10.5, ry+0.06, 2.3,  0.42, size=9.5, color=RPY_GRAY)

footer_strip(slide, "Lending (Razorpay Capital) sits at the top of the stack and generates the highest margin. "
             "Transactions at the bottom are a customer acquisition cost — not the core business.",
             highlight=True)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 9 — FEATURE ANALYSIS
# ════════════════════════════════════════════════════════════════════════
slide = new_slide("FEATURE ANALYSIS", "What's working, what's not, and why")
add_text(slide, "Feature Analysis", 0.6, 0.9, 12, 0.65, size=32, bold=True)

features = [
    ("Payment Gateway",   "✅ Strong",    "85%+ success rate; best-in-class DX; 8M+ merchants",
     "Commoditizing — PayU and Cashfree have closed the API gap"),
    ("RazorpayX",        "✅ Strategic", "API banking fills a genuine SME gap; 360° merchant view",
     "No banking licence — ceiling set by RBI's partner-bank model"),
    ("Razorpay Capital",  "✅ High potential", "Transaction-history underwriting beats bureau for SMEs",
     "Credit risk in an economic downturn; needs NPAs below 2%"),
    ("Opfin (Payroll)",  "⚠️ Execution risk", "Payroll data = best underwriting signal for Capital",
     "Compliance payroll (PF, ESI, TDS) is brutally hard to get right at scale"),
    ("Payment Links/Pages","✅ Strong",  "Democratized collection for non-dev merchants; Tier 2/3 reach",
     "No clear data pipeline from no-code users to Capital underwriting yet"),
    ("International",    "⚠️ Early",    "Addresses real pain for SaaS & exporters",
     "Established players (Wise, Payoneer); RBI cross-border rules are complex"),
]
col_x = [0.35, 2.55, 3.9, 7.9]
col_w = [2.1,  1.3,  3.9, 4.1]
row_y = 1.85
add_rect(slide, 0.35, row_y, 12.6, 0.36, fill_rgb=RPY_BLUE)
for hdr, cx, cw in zip(["Feature", "Signal", "What Works", "Risk / Gap"], col_x, col_w):
    add_text(slide, hdr, cx, row_y+0.04, cw, 0.3, size=11, bold=True, color=RPY_DARK)
for i, (feat, sig, works, risk) in enumerate(features):
    ry = row_y + 0.38 + i * 0.73
    bg = RPY_CARD if i % 2 == 0 else RGBColor(0x01, 0x1A, 0x38)
    add_rect(slide, 0.35, ry, 12.6, 0.71, fill_rgb=bg)
    add_text(slide, feat,  col_x[0], ry+0.08, col_w[0], 0.58, size=11, bold=True)
    sc = ACCENT_GREEN if "✅" in sig else ACCENT_AMBER
    add_text(slide, sig,   col_x[1], ry+0.08, col_w[1]+0.3, 0.58, size=10, color=sc)
    add_text(slide, works, col_x[2], ry+0.08, col_w[2], 0.58, size=10)
    add_text(slide, risk,  col_x[3], ry+0.08, col_w[3], 0.58, size=10, color=RPY_GRAY)

add_rect(slide, 0.35, 6.68, 12.6, 0.68, fill_rgb=RGBColor(0x01, 0x1A, 0x38))
for j, (label, val) in enumerate([
    ("Revenue driver:", "Gateway MDR + Capital interest income"),
    ("Biggest moat:", "Transaction data → SME credit underwriting"),
    ("Execution risk:", "Opfin compliance accuracy at scale"),
]):
    sx = 0.55 + j * 4.2
    add_rect(slide, sx, 6.71, 0.05, 0.62, fill_rgb=RPY_BLUE)
    add_text(slide, label, sx+0.14, 6.73, 1.5, 0.24, size=9, bold=True, color=RPY_BLUE)
    add_text(slide, val,   sx+0.14, 6.98, 3.9, 0.27, size=9.5, color=RPY_WHITE)


# ════════════════════════════════════════════════════════════════════════
# SECTION DIVIDER — THE BUSINESS
# ════════════════════════════════════════════════════════════════════════
section_divider("The Business", "Model · Revenue · Market Opportunity")


# ════════════════════════════════════════════════════════════════════════
# SLIDE 11 — BUSINESS MODEL
# ════════════════════════════════════════════════════════════════════════
slide = new_slide("BUSINESS MODEL", "The 3-layer stack — and why transactions are a CAC, not a business")
add_text(slide, "Business Model — 3 Layers", 0.6, 0.9, 12, 0.65, size=32, bold=True)

pillars = [
    ("01", "Transactions",
     "Payment gateway, Links, Pages, Smart Collect.\n\n"
     "Revenue: MDR ~1.8–3% per transaction.\n\n"
     "Strategic role: Customer acquisition cost. "
     "Thin margin; volume-dependent. The entry point "
     "into the relationship, not the revenue goal.",
     "Low–medium margin · Volume play"),
    ("02", "Banking",
     "RazorpayX, corporate cards, Opfin payroll.\n\n"
     "Revenue: Interchange, float income, SaaS fees.\n\n"
     "Strategic role: Deepens merchant lock-in; "
     "provides 360° cash flow view that enables Capital "
     "underwriting. Improves with AUM.",
     "Medium margin · Improves with AUM"),
    ("03", "Lending",
     "Razorpay Capital — business loans & credit lines.\n\n"
     "Revenue: Interest income + processing fees.\n\n"
     "Strategic role: Highest-margin business. "
     "Transaction-history underwriting = structurally "
     "better than bureau-based lending for SMEs. "
     "Credit risk is the constraint.",
     "High margin · Data flywheel advantage"),
]
for i, (num, title, desc, margin) in enumerate(pillars):
    x = 0.45 + i * 4.28
    add_rect(slide, x, 1.85, 4.08, 4.62, fill_rgb=RPY_CARD)
    add_rect(slide, x, 1.85, 4.08, 0.05, fill_rgb=RPY_BLUE)
    add_text(slide, num,   x+0.2, 1.94, 0.8, 0.55, size=26, bold=True, color=RPY_BLUE)
    add_text(slide, title, x+0.2, 2.55, 3.7, 0.52, size=15, bold=True)
    add_text(slide, desc,  x+0.2, 3.15, 3.7, 2.78, size=11.5, color=RPY_GRAY)
    add_rect(slide, x+0.2, 6.0,  3.6, 0.04, fill_rgb=RPY_BLUE)
    add_text(slide, margin, x+0.2, 6.1, 3.6, 0.3, size=10.5, bold=True, color=RPY_BLUE)

footer_strip(slide, "The long-term bet: gateway processing fees are a CAC. "
             "The real business is lending and banking — which require the transaction data that the gateway generates.",
             highlight=True)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 12 — BY THE NUMBERS
# ════════════════════════════════════════════════════════════════════════
slide = new_slide("BY THE NUMBERS", "Financial data FY21–FY24  ·  Sources: Inc42, Entrackr, public filings")
add_text(slide, "By the Numbers", 0.6, 0.9, 12, 0.65, size=32, bold=True)
add_img(slide, REV_CHART, 0.4, 1.68, 6.5, 2.85)

chart_insights = [
    (RPY_BLUE,   "~6.7× revenue growth", "₹390 Cr (FY21) → ₹2,600 Cr+ (FY24)"),
    (ACCENT_AMBER,"Net loss narrowing",   "Path to profitability via Capital margin expansion"),
    (ACCENT_GREEN,"Data flywheel effect", "More merchants → more transaction data → better underwriting → more Capital disbursals"),
]
for j, (col, heading, detail) in enumerate(chart_insights):
    cy = 4.7 + j * 0.82
    add_rect(slide, 0.4, cy, 6.5, 0.76, fill_rgb=RPY_CARD)
    add_rect(slide, 0.4, cy, 0.05, 0.76, fill_rgb=col)
    add_text(slide, heading, 0.58, cy+0.08, 5.7, 0.26, size=11, bold=True, color=col)
    add_text(slide, detail,  0.58, cy+0.38, 5.7, 0.32, size=10.5, color=RPY_GRAY)

stats = [
    ("8M+",   "Merchants on platform",     "As of FY24"),
    ("$90B+", "Annual payments processed", "GMV FY24"),
    ("85%+",  "Checkout success rate",     "vs 65–70% industry avg"),
    ("$7.5B", "Valuation — Series F",      "October 2021"),
    ("100+",  "Payment methods",           "UPI, cards, netbanking, wallets"),
    ("<2%",   "Capital NPA target",        "Transaction-history underwriting advantage"),
]
for i, (num, label, note) in enumerate(stats):
    row, col = i//3, i%3
    x = 7.2 + col * 2.05
    y = 1.68 + row * 2.45
    add_rect(slide, x, y, 1.95, 2.2, fill_rgb=RPY_CARD)
    add_rect(slide, x, y, 1.95, 0.05, fill_rgb=RPY_BLUE)
    add_text(slide, num,   x+0.12, y+0.18, 1.7, 0.7, size=20, bold=True, color=RPY_BLUE)
    add_text(slide, label, x+0.12, y+0.9,  1.7, 0.6, size=10, bold=True)
    add_text(slide, note,  x+0.12, y+1.52, 1.7, 0.62, size=9, color=RPY_GRAY)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 13 — HONEST ASSESSMENT
# ════════════════════════════════════════════════════════════════════════
slide = new_slide("HONEST ASSESSMENT", "What's working and what isn't")
add_text(slide, "Honest Assessment", 0.6, 0.9, 12, 0.65, size=32, bold=True)

add_rect(slide, 0.4, 1.85, 6.1, 0.42, fill_rgb=ACCENT_GREEN)
add_text(slide, "  ✓  What's Working", 0.4, 1.85, 6.1, 0.42, size=13.5, bold=True, color=RPY_DARK)
working = [
    "Developer experience moat — 8 years of compound advantage in API quality & docs",
    "Data flywheel — more transactions → better fraud detection → higher success rate → more merchants",
    "SME trust — 8M+ businesses rely on Razorpay for revenue collection; high switching cost",
    "Capital's underwriting advantage — transaction-history beats bureau for new-to-credit SMEs",
    "Enterprise expansion — moving upmarket without losing the SME base",
]
add_ml(slide, [("✓  " + w, False, RPY_WHITE) for w in working],
       0.4, 2.35, 5.9, 4.1, size=12.5)

add_rect(slide, 7.0, 1.85, 6.1, 0.42, fill_rgb=ACCENT_RED)
add_text(slide, "  ✗  What's Not Working", 7.0, 1.85, 6.1, 0.42, size=13.5, bold=True)
risks = [
    "Regulatory ceiling — no banking licence; one RBI circular away from a product constraint",
    "Commoditization in the gateway — PayU and Cashfree have closed the API gap",
    "Enterprise sales motion — self-serve breaks down above ₹10Cr transaction volume",
    "Opfin execution risk — payroll compliance errors at scale = reputational crisis",
    "Valuation overhang — $7.5B Series F set at 2021 peak; secondary market is difficult",
]
add_ml(slide, [("✗  " + r, False, RPY_WHITE) for r in risks],
       7.0, 2.35, 5.9, 4.1, size=12.5)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 14 — INDIA MARKET OPPORTUNITY
# ════════════════════════════════════════════════════════════════════════
slide = new_slide("MARKET OPPORTUNITY", "India's digital payments and SME finance landscape")
add_text(slide, "India Market Opportunity", 0.6, 0.9, 12, 0.65, size=32, bold=True)
add_img(slide, MARKET_MAP, 0.4, 1.68, 5.5, 3.5)

opp = [
    ("64M+",  "SMEs in India — Razorpay's addressable market"),
    ("8M+",   "Merchants on platform — ~12.5% TAM penetration"),
    ("₹12T+", "Annual digital payments TAM (FY24)"),
    ("$90B+", "Razorpay GMV — growing at ~30% YoY"),
    ("90%+",  "SMEs underserved by traditional finance → Capital TAM"),
    ("10×",   "Revenue/merchant headroom vs global benchmarks"),
]
add_rect(slide, 6.3, 1.7, 6.6, 0.4, fill_rgb=RPY_BLUE)
add_text(slide, "  Key Market Numbers", 6.3, 1.7, 6.6, 0.4, size=11.5, bold=True, color=RPY_DARK)
for i, (num, label) in enumerate(opp):
    ry = 2.14 + i * 0.6
    bg = RPY_CARD if i % 2 == 0 else RGBColor(0x01, 0x1A, 0x38)
    add_rect(slide, 6.3, ry, 6.6, 0.57, fill_rgb=bg)
    add_text(slide, num,   6.48, ry+0.08, 1.35, 0.35, size=14, bold=True, color=RPY_BLUE)
    add_text(slide, label, 7.88, ry+0.1,  4.7,  0.38, size=11)

footer_strip(slide, "Razorpay's opportunity is not user growth — it's revenue depth per merchant.  "
             "At current penetration levels, Capital alone has a multi-billion dollar runway within the existing merchant base.",
             highlight=True)


# ════════════════════════════════════════════════════════════════════════
# SECTION DIVIDER — THE COMPETITION
# ════════════════════════════════════════════════════════════════════════
section_divider("The Competition", "Razorpay vs Stripe · PayU · Cashfree · PhonePe Business")


# ════════════════════════════════════════════════════════════════════════
# SLIDE 16 — RAZORPAY vs STRIPE
# ════════════════════════════════════════════════════════════════════════
slide = new_slide("RAZORPAY vs STRIPE", "Why the 'Indian Stripe' framing misses the point")
add_text(slide, "Razorpay vs Stripe — Key Differences",
         0.6, 0.9, 12, 0.65, size=30, bold=True)

comp_rows = [
    ("Dimension",          "Razorpay",                      "Stripe"),
    ("Core user",          "Indian developer + SME finance", "Global developer, internet businesses"),
    ("Market context",     "UPI-first; large unbanked SME",  "Card-first; established banking infra"),
    ("Revenue emphasis",   "Transactions + banking + lending","Transactions + Stripe Capital"),
    ("Data advantage",     "Deep India SME cash flow data",  "Global transaction patterns"),
    ("Regulatory posture", "Navigates RBI, SEBI, IRDAI",    "Works within established frameworks"),
    ("Key moat",           "India-specific data + compliance","Global developer network effects"),
    ("Expansion limit",    "India-specific advantages don't\ntravel well",
                          "Global — but India requires local\nregulatory relationships"),
]
col_x = [0.35, 4.15, 8.9]; col_w = [3.7, 4.7, 4.35]
for i, row in enumerate(comp_rows):
    ry = 1.78 + i * 0.64
    if i == 0:
        add_rect(slide, 0.35, ry, 12.6, 0.64, fill_rgb=RPY_BLUE)
        tc, bld = RPY_DARK, True
    else:
        bg = RPY_CARD if i % 2 == 1 else RGBColor(0x01, 0x1A, 0x38)
        add_rect(slide, 0.35, ry, 12.6, 0.64, fill_rgb=bg)
        tc, bld = RPY_WHITE, False
    for j, (cell, cx, cw) in enumerate(zip(row, col_x, col_w)):
        if j == 1 and i > 0:
            add_rect(slide, cx-0.05, ry, cw+0.1, 0.64,
                     fill_rgb=RGBColor(0x01, 0x22, 0x48))
            cc = RPY_BLUE
        else:
            cc = tc if i > 0 else RPY_DARK
        add_text(slide, cell, cx, ry+0.08, cw, 0.52,
                 size=10.5, bold=(bld or j == 0), color=cc)

footer_strip(slide, "Key insight: Razorpay can build products Stripe cannot — India-specific data, "
             "regulatory relationships, and SME distribution. The moat protects them in India; expansion carries none of these advantages.",
             highlight=True)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 17 — COMPETITIVE LANDSCAPE
# ════════════════════════════════════════════════════════════════════════
slide = new_slide("COMPETITIVE LANDSCAPE 2025–26", "The real challengers to Razorpay's gateway moat")
add_text(slide, "Indian Fintech Competitive Landscape", 0.6, 0.9, 12, 0.65, size=30, bold=True)

comp_data = [
    (LOGO_PAYU,     "PayU India",    "Enterprises,\nmid-market",
     "Prosus-backed,\nlong-standing PG",  "International,\nemi options",  "HIGH",   "Direct gateway competitor; comparable API now"),
    (LOGO_CASHFREE, "Cashfree",      "Startups,\nfintechs",
     "Strong payouts\n& disbursals",      "Faster payouts\nthan Razorpay",     "HIGH",   "Developer trust; strong in payout use-cases"),
    (LOGO_JUSPAY,   "Juspay",        "Large enterprises\n(Swiggy, CRED)",
     "Hyperswitch\nmulti-PSP router",     "Checkout\northestration",       "MED-HIGH","Addresses checkout failure at scale — PM layer"),
    (LOGO_PHONEPE,  "PhonePe Biz",   "SMBs,\noffline merchants",
     "UPI + SmartSpeaker\nedge", "Zero MDR UPI,\nkirana reach",          "MEDIUM", "Bharat segment; different buyer persona"),
    (LOGO_STRIPE,   "Stripe",        "Global SaaS,\nstartups",
     "Global infra,\nradar fraud",        "Best-in-class DX\nfor global cos",  "LOW",    "No Indian-specific compliance depth"),
    (LOGO_CCAVENUE, "CCAvenue",      "SMEs, legacy\nenterprises",
     "20+ years,\nbank relationships",    "High-value\nenterprise legacy",    "LOW",    "API quality gap; losing developer mindshare"),
]
col_xs = [0.35, 1.55, 3.1, 5.1, 7.05, 8.95, 10.95]
col_ws = [1.05, 1.5,  1.9, 1.9, 1.85, 1.95, 2.25]
hdrs   = ["",  "Player", "Target",  "Edge",   "Differentiator",  "Threat",  "PM Note"]

add_rect(slide, 0.35, 1.78, 12.6, 0.36, fill_rgb=RPY_BLUE)
for hdr, cx, cw in zip(hdrs, col_xs, col_ws):
    add_text(slide, hdr, cx, 1.8, cw, 0.3, size=10, bold=True, color=RPY_DARK)

for i, (logo_path, name, target, edge, diff, threat, note) in enumerate(comp_data):
    ry = 2.18 + i * 0.82
    bg = RPY_CARD if i % 2 == 0 else RGBColor(0x01, 0x1A, 0x38)
    add_rect(slide, 0.35, ry, 12.6, 0.8, fill_rgb=bg)
    add_img(slide, logo_path, col_xs[0]+0.05, ry+0.1, 0.6, 0.6)
    tc = ACCENT_RED if threat == "HIGH" else (ACCENT_AMBER if "MED" in threat else ACCENT_GREEN)
    for j, (val, cx, cw) in enumerate(zip([name, target, edge, diff, threat, note],
                                           col_xs[1:], col_ws[1:])):
        c = tc if j == 4 else RPY_WHITE
        add_text(slide, val, cx, ry+0.08, cw, 0.68, size=10,
                 bold=(j == 0 or j == 4), color=c)


# ════════════════════════════════════════════════════════════════════════
# SECTION DIVIDER — PRODUCT THINKING
# ════════════════════════════════════════════════════════════════════════
section_divider("Product Thinking", "3 Opportunities I'd Build Next · PM Takeaways")


# ════════════════════════════════════════════════════════════════════════
# SLIDE 19 — PRODUCT OPPORTUNITIES
# ════════════════════════════════════════════════════════════════════════
slide = new_slide("PRODUCT OPPORTUNITIES", "What I would build next — and why now")
add_text(slide, "Three Product Opportunities", 0.6, 0.9, 12, 0.65, size=32, bold=True)

opps = [
    ("01", "AI Cash Flow\nIntelligence",
     "Problem: Razorpay Capital is reactive — merchants apply when they need money.\n\n"
     "Build: Forward-looking cash flow forecast using transaction time-series. "
     "'You have 23 days of runway. Your August dip is coming — here's a pre-approved offer.'\n\n"
     "Why: Proactive credit offer → 3–5× conversion vs. cold pitch. The data already exists.",
     "RICE ~580 · HIGH"),
    ("02", "Razorpay for Bharat",
     "Problem: 8M merchants ≠ 64M SMEs. The remaining 56M are Tier 2/3 businesses "
     "on WhatsApp, in Hindi/Marathi/Tamil, with unreliable internet.\n\n"
     "Build: Simplified vernacular-first, low-bandwidth payment collection. "
     "Offline-first. Feature phone compatible.\n\n"
     "Why: UPI solved the consumer side. Merchant-side collection infrastructure for Bharat is unbuilt.",
     "RICE ~430 · HIGH"),
    ("03", "Embedded\nCompliance",
     "Problem: SMEs spend 10–15 hrs/month on GST filings, TDS, PF — or pay ₹3,000+/month to CAs.\n\n"
     "Build: Pre-populate GST/TDS filings from transaction data. Review + file in one click.\n\n"
     "Why: Makes switching from Razorpay almost impossible. Compliance history lives in the platform. "
     "A local moat Stripe cannot replicate.",
     "RICE ~350 · HIGH"),
]
for i, (num, title, desc, rice) in enumerate(opps):
    x = 0.4 + i * 4.28
    add_rect(slide, x, 1.85, 4.1, 4.98, fill_rgb=RPY_CARD)
    add_rect(slide, x, 1.85, 4.1, 0.05, fill_rgb=RPY_BLUE)
    add_text(slide, num,   x+0.2, 1.94, 0.8, 0.5,  size=22, bold=True, color=RPY_BLUE)
    add_text(slide, title, x+0.2, 2.5,  3.7, 0.72,  size=14, bold=True)
    add_text(slide, desc,  x+0.2, 3.28, 3.7, 3.18, size=11, color=RPY_GRAY)
    add_rect(slide, x+0.2, 6.4,  3.6, 0.3, fill_rgb=RGBColor(0x01, 0x1A, 0x38))
    add_text(slide, rice,  x+0.3, 6.43, 3.4, 0.26, size=10.5, bold=True, color=RPY_BLUE)

add_rect(slide, 0.4, 6.82, 12.5, 0.45, fill_rgb=RGBColor(0x01, 0x1A, 0x38))
add_text(slide, "Recommended sequence:  AI Cash Flow → Embedded Compliance → Razorpay for Bharat",
         0.6, 6.88, 12.2, 0.36, size=11, color=RPY_BLUE)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 20 — PM TAKEAWAYS
# ════════════════════════════════════════════════════════════════════════
slide = new_slide("PM TAKEAWAYS", "What this teardown taught me")
add_text(slide, "PM Takeaways", 0.6, 0.9, 12, 0.65, size=32, bold=True)

takeaways = [
    ("01", "Developer experience is a product, not a feature.",
     "Razorpay won its first market because documentation and API ergonomics were better — "
     "not because success rates were higher. In developer-bought products, the integration experience "
     "is the first product the buyer evaluates."),
    ("02", "Platform expansions should use each other's data.",
     "Every Razorpay product addition was justified partly because it generated data that improved an "
     "existing product. Adjacency without data leverage is just scope creep."),
    ("03", "In regulated fintech, compliance is both constraint and moat.",
     "Embedded compliance exists only because Razorpay has the regulatory relationships to file on behalf "
     "of merchants. Competitors can't copy the feature without copying the regulatory foundation. "
     "Build compliance in — don't bolt it on."),
    ("04", "Transaction revenue is a customer acquisition cost in disguise.",
     "The long-term bet is that payment processing margins pay for merchant acquisition, and the real "
     "business is lending and banking. Understanding multi-product LTV changes which merchants are "
     "worth acquiring at what cost."),
    ("05", "Design for the skeptic, not the enthusiast.",
     "Razorpay's finance-team-friendly dashboard wasn't technically necessary — developers could read "
     "API logs. But the finance team has veto power. Designing for the person who can say no is "
     "as important as designing for the one who wants to use it."),
]
for i, (num, title, desc) in enumerate(takeaways):
    ry = 1.82 + i * 1.1
    add_rect(slide, 0.4, ry, 12.5, 1.04, fill_rgb=RPY_CARD)
    add_rect(slide, 0.4, ry, 0.06, 1.04, fill_rgb=RPY_BLUE)
    add_text(slide, num,   0.62, ry+0.1, 0.6,  0.42, size=14, bold=True, color=RPY_BLUE)
    add_text(slide, title, 1.28, ry+0.08, 10.0, 0.36, size=13, bold=True)
    add_text(slide, desc,  1.28, ry+0.46, 11.4, 0.55, size=11, color=RPY_GRAY)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 21 — CLOSING
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=RGBColor(0x00, 0x0C, 0x20))
add_rect(slide, 0, 0, 13.33, 0.07, fill_rgb=RPY_BLUE)
add_rect(slide, 0, 7.43, 13.33, 0.07, fill_rgb=RPY_BLUE)
add_rect(slide, 0, 0, 0.14, 7.5, fill_rgb=RPY_BLUE)
add_rect(slide, 6.5, 0, 6.83, 7.5, fill_rgb=RPY_NAVY)

add_text(slide, "Shraddha Singh",
         0.7, 1.6, 7, 0.8, size=36, bold=True)
add_text(slide, "Product Manager",
         0.7, 2.55, 6, 0.48, size=18, color=RPY_BLUE)
add_text(slide, "6 years · Fintech · Insurtech · Enterprise GenAI\n"
         "KYC onboarding · Healthcare AI · Insurance platform",
         0.7, 3.12, 6.5, 0.7, size=13, color=RPY_GRAY)
add_rect(slide, 0.7, 3.95, 2.8, 0.04, fill_rgb=RPY_BLUE)
add_text(slide, "Open to PM · TPM · AI PM roles",
         0.7, 4.1, 6, 0.38, size=13, bold=True)
add_text(slide, "shraddhasingh8968@gmail.com",
         0.7, 4.62, 6.5, 0.35, size=12, color=RPY_BLUE)
add_text(slide, "linkedin.com/in/shraddha-singh-9149a4172",
         0.7, 5.0, 6.5, 0.35, size=12, color=RPY_GRAY)
add_text(slide, "shraddhasingh.dev",
         0.7, 5.38, 4, 0.35, size=12, color=RPY_GRAY)

add_img(slide, LOGO_RAZORPAY if os.path.exists(str(LOGO_RAZORPAY)) else RPY_LOGO, 7.5, 2.2, 4.8, 1.6)
add_text(slide, "Product Teardown · June 2026",
         7.5, 4.0, 5.5, 0.42, size=13, color=RPY_GRAY, align=PP_ALIGN.CENTER)
add_text(slide, "Razorpay — From Payment Gateway to Financial OS",
         7.3, 4.55, 5.8, 0.55, size=12.5, italic=True, color=RPY_LGRAY, align=PP_ALIGN.CENTER)

# ── Save ─────────────────────────────────────────────────────────────────
out_path = os.path.join(
    os.path.dirname(os.path.abspath(__file__)),
    "Razorpay-Teardown-Shraddha-Singh.pptx"
)
prs.save(out_path)
print(f"\nSaved: {out_path}")
print(f"Slides: {len(prs.slides)}")
