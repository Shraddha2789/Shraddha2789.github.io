"""
CRED Product Teardown — Complete 25-Slide Deck
Shraddha Singh · Feb 2026
"""

import os
from io import BytesIO
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colors ──────────────────────────────────────────────────────
CRED_BLACK   = RGBColor(0x1A, 0x1A, 0x1A)
CRED_GOLD    = RGBColor(0xC9, 0xA8, 0x4C)
CRED_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
CRED_GRAY    = RGBColor(0xB0, 0xB0, 0xB0)
CRED_DARK    = RGBColor(0x12, 0x12, 0x12)
ACCENT_RED   = RGBColor(0xE5, 0x3E, 0x3E)
ACCENT_GREEN = RGBColor(0x3E, 0xC9, 0x7A)
ACCENT_AMBER = RGBColor(0xF5, 0xA6, 0x23)
ACCENT_BLUE  = RGBColor(0x4A, 0x90, 0xD9)

ASSETS = "/tmp/cred_assets"
os.makedirs(ASSETS, exist_ok=True)


# ════════════════════════════════════════════════════════════════
# IMAGE GENERATION HELPERS (Pillow)
# ════════════════════════════════════════════════════════════════

def make_logo_badge(text, bg_hex, fg_hex, size=200, radius=28, filepath=None):
    """Create a rounded-rect logo badge with text."""
    img = Image.new("RGBA", (size, size), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)
    bg = tuple(int(bg_hex[i:i+2], 16) for i in (0, 2, 4)) + (255,)
    fg = tuple(int(fg_hex[i:i+2], 16) for i in (0, 2, 4)) + (255,)
    draw.rounded_rectangle([0, 0, size-1, size-1], radius=radius, fill=bg)
    try:
        font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", size=size//4)
    except:
        font = ImageFont.load_default()
    bbox = draw.textbbox((0, 0), text, font=font)
    tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
    draw.text(((size - tw) // 2, (size - th) // 2 - 4), text, font=font, fill=fg)
    path = filepath or f"{ASSETS}/{text.lower().replace(' ','_')}.png"
    img.save(path)
    return path


def make_cred_logo(filepath=None):
    """CRED wordmark — black bg, gold lettering."""
    w, h = 480, 180
    img = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)
    draw.rounded_rectangle([0, 0, w-1, h-1], radius=20, fill=(26, 26, 26, 255))
    draw.rounded_rectangle([0, 0, w-1, 8], radius=4, fill=(201, 168, 76, 255))
    try:
        font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", size=90)
    except:
        font = ImageFont.load_default()
    bbox = draw.textbbox((0, 0), "CRED", font=font)
    tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
    draw.text(((w - tw) // 2, (h - th) // 2 + 4), "CRED", font=font,
              fill=(201, 168, 76, 255))
    path = filepath or f"{ASSETS}/cred_logo.png"
    img.save(path)
    return path


def make_phone_mockup(title, lines, accent_hex="C9A84C", filepath=None):
    """Draw a phone-frame mockup with app content."""
    W, H = 320, 580
    img = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)
    # phone body
    draw.rounded_rectangle([4, 4, W-5, H-5], radius=38,
                            fill=(22, 22, 22, 255), outline=(80, 80, 80, 255), width=2)
    # status bar
    draw.rectangle([24, 24, W-24, 46], fill=(30, 30, 30, 255))
    draw.ellipse([W//2-14, 16, W//2+14, 32], fill=(20, 20, 20, 255))  # notch
    # accent header bar
    accent = tuple(int(accent_hex[i:i+2], 16) for i in (0, 2, 4)) + (255,)
    draw.rectangle([24, 50, W-24, 100], fill=accent)
    try:
        font_lg = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", size=20)
        font_sm = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", size=14)
        font_xs = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", size=12)
    except:
        font_lg = font_sm = font_xs = ImageFont.load_default()
    draw.text((40, 62), title, font=font_lg, fill=(26, 26, 26, 255))
    # content rows
    y = 115
    for i, (label, val, bar_ratio) in enumerate(lines):
        row_bg = (28, 28, 28, 255) if i % 2 == 0 else (32, 32, 32, 255)
        draw.rectangle([24, y, W-24, y+48], fill=row_bg)
        draw.text((36, y+8), label, font=font_xs, fill=(176, 176, 176, 255))
        draw.text((36, y+26), val, font=font_sm, fill=(255, 255, 255, 255))
        if bar_ratio > 0:
            bar_w = int((W-80) * bar_ratio)
            draw.rounded_rectangle([W-50, y+18, W-50+bar_w, y+30],
                                   radius=4, fill=accent)
        y += 52
    # home indicator
    draw.rounded_rectangle([W//2-30, H-22, W//2+30, H-14],
                           radius=4, fill=(80, 80, 80, 255))
    path = filepath or f"{ASSETS}/phone_{title.replace(' ','_')}.png"
    img.save(path)
    return path


def make_bar_chart(data, title, filepath=None):
    """Simple horizontal bar chart image."""
    W, H = 700, 320
    img = Image.new("RGBA", (W, H), (18, 18, 18, 255))
    draw = ImageDraw.Draw(img)
    try:
        font_hd = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", size=20)
        font_sm = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", size=16)
    except:
        font_hd = font_sm = ImageFont.load_default()
    draw.text((20, 16), title, font=font_hd, fill=(201, 168, 76, 255))
    bar_x = 140; max_bar = W - bar_x - 120
    bar_h = 42; gap = 22
    max_val = max(v for _, v, _ in data)
    gold  = (201, 168, 76, 255)
    gray  = (50, 50, 50, 255)
    white = (255, 255, 255, 255)
    lt    = (176, 176, 176, 255)
    y = 60
    for label, val, color_hex in data:
        ratio = val / max_val
        bar_color = tuple(int(color_hex[i:i+2], 16) for i in (0,2,4)) + (255,)
        draw.rectangle([bar_x, y, bar_x + int(max_bar*ratio), y+bar_h],
                       fill=bar_color)
        draw.text((14, y + 10), label, font=font_sm, fill=white)
        draw.text((bar_x + int(max_bar*ratio) + 10, y+10),
                  f"₹{val:,} Cr", font=font_sm, fill=gold)
        y += bar_h + gap
    path = filepath or f"{ASSETS}/chart_{title.replace(' ','_')}.png"
    img.save(path)
    return path


def make_india_map_placeholder(filepath=None):
    """Styled India market size visual."""
    W, H = 600, 380
    img = Image.new("RGBA", (W, H), (18, 18, 18, 255))
    draw = ImageDraw.Draw(img)
    try:
        font_xl = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", size=52)
        font_lg = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", size=28)
        font_sm = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", size=18)
    except:
        font_xl = font_lg = font_sm = ImageFont.load_default()
    gold = (201, 168, 76, 255)
    gray = (176, 176, 176, 255)
    white = (255, 255, 255, 255)
    # concentric circles as market funnel
    draw.ellipse([160, 40, 440, 340], outline=(50, 50, 50, 255), width=2)
    draw.ellipse([200, 80, 400, 300], outline=(80, 80, 80, 255), width=2)
    draw.ellipse([240, 120, 360, 260], fill=(40, 34, 16, 255),
                 outline=(201, 168, 76, 255), width=2)
    draw.text((268, 168), "13M", font=font_lg, fill=gold)
    draw.text((250, 202), "CRED MAU", font=font_sm, fill=gray)
    draw.text((32, 60), "100M+\nCredit\nCard\nHolders", font=font_sm, fill=gray)
    draw.text((460, 60), "35–40M\nPremium\n(750+ CIBIL)", font=font_sm, fill=gray)
    draw.text((32, 290), "India TAM", font=font_sm, fill=gray)
    draw.line([120, 155, 230, 190], fill=gray, width=1)
    draw.line([460, 155, 380, 190], fill=gold, width=1)
    path = filepath or f"{ASSETS}/india_market.png"
    img.save(path)
    return path


# ── Generate all assets ─────────────────────────────────────────
print("Generating image assets...")
CRED_LOGO   = make_cred_logo()
LOGO_OC     = make_logo_badge("OC",     "1A1A2E", "C9A84C", filepath=f"{ASSETS}/onecard.png")
LOGO_SLICE  = make_logo_badge("slice",  "FF4D4D", "FFFFFF", filepath=f"{ASSETS}/slice.png")
LOGO_UNI    = make_logo_badge("uni",    "2D6A4F", "FFFFFF", filepath=f"{ASSETS}/uni.png")
LOGO_JUP    = make_logo_badge("JM",     "5B2D8E", "FFFFFF", filepath=f"{ASSETS}/jupiter.png")
LOGO_FI     = make_logo_badge("fi",     "00C896", "1A1A1A", filepath=f"{ASSETS}/fi.png")
LOGO_SCAPIA   = make_logo_badge("SC",    "0077B6", "FFFFFF", filepath=f"{ASSETS}/scapia.png")
LOGO_PHONEPE  = make_logo_badge("PP",   "5F259F", "FFFFFF", filepath=f"{ASSETS}/phonepe.png")
LOGO_PAYTM    = make_logo_badge("PTM",  "00BAF2", "FFFFFF", filepath=f"{ASSETS}/paytm.png")
LOGO_GPAY     = make_logo_badge("G",    "4285F4", "FFFFFF", filepath=f"{ASSETS}/gpay.png")

PHONE_BILLS = make_phone_mockup("My Bills", [
    ("HDFC Regalia",   "₹18,450 due",  0.72),
    ("ICICI Sapphire", "₹6,220 due",   0.25),
    ("Axis Magnus",    "₹32,100 due",  1.00),
    ("Coins earned",   "2,847 CRED ◆", 0.0),
    ("Next payment",   "Mar 5, 2026",  0.0),
], accent_hex="C9A84C")

PHONE_CASH = make_phone_mockup("CRED Cash", [
    ("Pre-approved limit", "₹4,00,000",  1.00),
    ("Interest rate",      "1.5% / mo",  0.30),
    ("Processing fee",     "Zero",       0.0),
    ("Disbursal",          "< 2 minutes",0.0),
    ("CIBIL impact",       "Soft pull",  0.0),
], accent_hex="3EC97A")

PHONE_GARAGE = make_phone_mockup("CRED Garage", [
    ("Honda City",     "MH02-AX-4521",  0.0),
    ("FASTag balance", "₹420 — Recharge",0.45),
    ("Insurance",      "Renew by Apr 3", 0.9),
    ("Last service",   "42 days ago",   0.0),
    ("Fuel savings",   "₹340 this month",0.0),
], accent_hex="4A90D9")

REV_CHART = make_bar_chart([
    ("FY 2022",  394,   "555555"),
    ("FY 2023",  1400,  "8B7332"),
    ("FY 2024",  2473,  "C9A84C"),
], "Operating Revenue (₹ Crore)")

INDIA_MKT = make_india_map_placeholder()
print("Assets ready.")


# ════════════════════════════════════════════════════════════════
# PPT HELPERS
# ════════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_pt=0):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.line.width = Pt(line_pt)
    if fill_rgb:
        s.fill.solid(); s.fill.fore_color.rgb = fill_rgb
    else:
        s.fill.background()
    if line_rgb:
        s.line.color.rgb = line_rgb
    else:
        s.line.fill.background()
    return s


def add_text(slide, text, l, t, w, h, size=14, bold=False, color=CRED_WHITE,
             align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tb


def add_ml(slide, lines, l, t, w, h, size=12, default_color=CRED_WHITE):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True; tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            txt, bld, clr = item, False, default_color
        else:
            txt = item[0]
            bld = item[1] if len(item) > 1 else False
            clr = item[2] if len(item) > 2 else default_color
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(3)
        r = p.add_run(); r.text = txt
        r.font.size = Pt(size); r.font.bold = bld; r.font.color.rgb = clr
    return tb


def add_img(slide, path, l, t, w, h=None):
    try:
        if h:
            slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
        else:
            slide.shapes.add_picture(path, Inches(l), Inches(t), width=Inches(w))
    except Exception as e:
        print(f"  img skip: {e}")


def dark_bg(slide, with_bars=True):
    add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=CRED_DARK)
    if with_bars:
        add_rect(slide, 0, 0,    13.33, 0.06, fill_rgb=CRED_GOLD)
        add_rect(slide, 0, 7.44, 13.33, 0.06, fill_rgb=CRED_GOLD)


def section_label(slide, tag, subtitle=""):
    add_text(slide, tag,      0.6, 0.18, 12, 0.38, size=10, color=CRED_GOLD)
    if subtitle:
        add_text(slide, subtitle, 0.6, 0.48, 12, 0.32, size=9, color=CRED_GRAY)


def new_slide(tag=None, subtitle=None):
    s = prs.slides.add_slide(BLANK)
    dark_bg(s)
    if tag:
        section_label(s, tag, subtitle)
    return s


def section_divider(label, subtitle=""):
    """Full-bleed dark section break with gold accent."""
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=RGBColor(0x10, 0x10, 0x10))
    add_rect(slide, 0, 0, 0.12, 7.5, fill_rgb=CRED_GOLD)
    add_rect(slide, 0, 0, 13.33, 0.05, fill_rgb=CRED_GOLD)
    add_rect(slide, 0, 7.45, 13.33, 0.05, fill_rgb=CRED_GOLD)
    add_text(slide, label,    0.5, 2.8, 12, 1.2, size=52, bold=True,
             color=CRED_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.5, 4.2, 12, 0.55, size=16,
                 color=CRED_GOLD, align=PP_ALIGN.LEFT)
    add_text(slide, "CRED  ·  Product Teardown  ·  Shraddha Singh  ·  Feb 2026",
             0.5, 7.0, 12, 0.38, size=10, color=CRED_GRAY)
    return slide


# ════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=RGBColor(0x0E, 0x0E, 0x0E))
add_rect(slide, 0, 0,    13.33, 0.08, fill_rgb=CRED_GOLD)
add_rect(slide, 0, 7.42, 13.33, 0.08, fill_rgb=CRED_GOLD)
# diagonal accent
add_rect(slide, 7.5, 0, 5.83, 7.5, fill_rgb=RGBColor(0x16, 0x14, 0x0A))
add_img(slide, CRED_LOGO, 7.8, 1.8, 4.5, 1.7)
add_text(slide, "CRED", 0.9, 1.2, 7, 2.4, size=108, bold=True, color=CRED_WHITE)
add_text(slide, "Product Teardown", 0.9, 3.85, 7, 0.7, size=30, color=CRED_GOLD)
add_text(slide, "An inside-out analysis of India's most opinionated fintech product",
         0.9, 4.65, 7, 0.55, size=14, color=CRED_GRAY)
add_rect(slide, 0.9, 5.5, 3, 0.04, fill_rgb=CRED_GOLD)
add_text(slide, "Shraddha Singh", 0.9, 5.65, 5, 0.38, size=14, bold=True)
add_text(slide, "Product Manager  ·  Fintech & Enterprise SaaS",
         0.9, 6.08, 6, 0.35, size=12, color=CRED_GRAY)
add_text(slide, "Feb 2026", 0.9, 6.5, 3, 0.32, size=11, color=CRED_GOLD)


# ════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ════════════════════════════════════════════════════════════════
slide = new_slide("AGENDA", "What this teardown covers")
add_text(slide, "Agenda", 0.6, 0.9, 12, 0.7, size=36, bold=True)

sections = [
    ("01", "Context",          "What is CRED · Target user · Core problem"),
    ("02", "The Product",      "Product evolution · Feature analysis · App UI showcase"),
    ("03", "The Business",     "Business model · Revenue data · Market opportunity"),
    ("04", "The Competition",  "4-player landscape · 6-way fintech threat analysis"),
    ("05", "Product Thinking", "Opportunities · User pain points · Onboarding · NSM"),
]
for i, (num, title, detail) in enumerate(sections):
    row, col = i // 3, i % 3
    if i < 3:
        x, y = 0.5 + col * 4.3, 1.85
    else:
        x, y = 0.5 + (i-3) * 6.5, 4.55

    w = 4.0 if i < 3 else 6.0
    add_rect(slide, x, y, w, 2.2, fill_rgb=RGBColor(0x22, 0x22, 0x22))
    add_rect(slide, x, y, w, 0.06, fill_rgb=CRED_GOLD)
    add_text(slide, num,    x+0.2, y+0.12, 0.8, 0.5, size=22, bold=True, color=CRED_GOLD)
    add_text(slide, title,  x+0.2, y+0.65, w-0.4, 0.5, size=15, bold=True)
    add_text(slide, detail, x+0.2, y+1.2,  w-0.4, 0.85, size=11, color=CRED_GRAY)


# ════════════════════════════════════════════════════════════════
# SLIDE 3 — WHAT IS CRED
# ════════════════════════════════════════════════════════════════
slide = new_slide("WHAT IS CRED", "Overview")
add_text(slide, "What is CRED?", 0.6, 0.9, 10, 0.7, size=36, bold=True)
add_rect(slide, 0.6, 1.75, 0.05, 0.9, fill_rgb=CRED_GOLD)
add_text(slide, "Pay your credit card bills. Get rewarded. Be part of an exclusive club.",
         0.85, 1.75, 8.5, 0.5, size=16, italic=True, color=CRED_GOLD)
bullets = [
    "Founded 2018 by Kunal Shah (FreeCharge)",
    "Invite-only: only users with CIBIL 750+ can join",
    "Core utility: aggregate all credit cards, pay bills in one tap",
    "Reward layer: earn CRED coins for every rupee paid on time",
    "Now expanding into lending, travel, investments, UPI & vehicle management",
]
add_ml(slide, [("• " + b, False, CRED_WHITE) for b in bullets],
       0.85, 2.4, 8.0, 2.8, size=14)   # reduced height to make room for KPI strip
add_img(slide, CRED_LOGO, 9.8, 1.85, 2.9, 1.1)

# KPI strip — 4 headline numbers
kpi_data = [("13M", "Monthly Active Users"), ("₹2,473 Cr", "FY24 Revenue"),
            ("750+", "Min. CIBIL to Join"), ("₹19K Cr", "Lending AUM")]
for i, (num, label) in enumerate(kpi_data):
    kx = 0.6 + i * 3.15
    add_rect(slide, kx, 5.38, 3.0, 0.92, fill_rgb=RGBColor(0x22, 0x22, 0x22))
    add_rect(slide, kx, 5.38, 3.0, 0.06, fill_rgb=CRED_GOLD)
    add_text(slide, num,   kx+0.18, 5.48, 2.7, 0.44, size=20, bold=True, color=CRED_GOLD)
    add_text(slide, label, kx+0.18, 5.94, 2.7, 0.3,  size=10, color=CRED_GRAY)

add_rect(slide, 0.6, 6.45, 12.1, 0.88, fill_rgb=RGBColor(0x26, 0x26, 0x26))
add_text(slide, "\"The credit card bill was never the product.  The user was.\"",
         0.9, 6.5, 11.8, 0.78, size=14, italic=True, color=CRED_GOLD, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 4 — THE USER
# ════════════════════════════════════════════════════════════════
slide = new_slide("THE USER", "Who CRED is built for")
add_text(slide, "The Target User", 0.6, 0.9, 10, 0.7, size=36, bold=True)
add_rect(slide, 0.4, 1.8, 5.9, 4.6, fill_rgb=RGBColor(0x22, 0x22, 0x22))
add_rect(slide, 0.4, 1.8, 5.9, 0.06, fill_rgb=CRED_GOLD)
add_text(slide, "Persona — Aspirational Aryan", 0.6, 1.9, 5.6, 0.5, size=13, bold=True, color=CRED_GOLD)
persona = ["Age: 26–40, urban professional", "Income: ₹10L+ per annum",
           "Owns 2–4 credit cards", "Values exclusivity, design, and his time",
           "Pays bills but gets zero back from his bank",
           "Trusts brands that treat him as premium"]
add_ml(slide, [("• " + p, False, CRED_WHITE) for p in persona], 0.6, 2.5, 5.5, 3.8, size=13)

add_rect(slide, 6.6, 1.8, 6.3, 4.6, fill_rgb=RGBColor(0x22, 0x22, 0x22))
add_rect(slide, 6.6, 1.8, 6.3, 0.06, fill_rgb=CRED_GOLD)
add_text(slide, "Why this user = CRED's product asset", 6.8, 1.9, 6.0, 0.5, size=13, bold=True, color=CRED_GOLD)
assets_data = [
    ("Low credit risk",       "→ Safe to build a lending book"),
    ("High purchasing power", "→ Valuable to brands & advertisers"),
    ("Multi-card holder",     "→ Rich spend data across categories"),
    ("Aspirational mindset",  "→ Responds to status-driven design"),
]
y = 2.55
for label, desc in assets_data:
    add_rect(slide, 6.6, y, 6.3, 0.65, fill_rgb=RGBColor(0x1E, 0x1E, 0x1E))
    add_text(slide, label, 6.8, y+0.08, 2.8, 0.4, size=13, bold=True, color=CRED_GOLD)
    add_text(slide, desc,  9.7, y+0.12, 3.0, 0.4, size=12)
    y += 0.72
# Summary insight for right column bottom — fills empty space
add_rect(slide, 6.6, 5.5, 6.3, 0.88, fill_rgb=RGBColor(0x25, 0x22, 0x12))
add_rect(slide, 6.6, 5.5, 0.08, 0.88, fill_rgb=CRED_GOLD)
add_text(slide, "Result: CRED can underwrite, advertise to, and retain this cohort "
         "at a fraction of normal CAC. The 750+ filter is not a constraint — it is a product decision.",
         6.82, 5.56, 5.92, 0.76, size=10, italic=True, color=CRED_GOLD)

add_rect(slide, 0.4, 6.55, 12.5, 0.06, fill_rgb=CRED_GOLD)
add_text(slide, "CRED didn't just pick a demographic. It picked a data asset.",
         0.4, 6.65, 12.5, 0.5, size=12, italic=True, color=CRED_GRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 5 — CORE PROBLEM
# ════════════════════════════════════════════════════════════════
slide = new_slide("CORE PROBLEM", "The insight behind CRED")
add_text(slide, "The Problem CRED Solves", 0.6, 0.9, 12, 0.7, size=36, bold=True)
problems = [
    ("Bill payment is boring & unrewarded",
     "Banks give nothing back for on-time payments — a 2/10 experience"),
    ("Premium buyers are expensive to reach",
     "Brands can't efficiently target creditworthy, high-income users at scale"),
    ("Good credit behaviour is unrecognised",
     "No social or economic reward for financial discipline in India"),
]
impact_tags = ["Bill-pay disruption", "Premium ad market ($B+ TAM)", "Credit culture reform"]
for i, (prob, insight) in enumerate(problems):
    x = 0.45 + i * 4.3
    add_rect(slide, x, 1.9, 4.1, 4.0, fill_rgb=RGBColor(0x22, 0x22, 0x22))
    add_rect(slide, x, 1.9, 4.1, 0.06, fill_rgb=CRED_GOLD)
    add_text(slide, f"0{i+1}", x+0.2, 2.02, 0.7, 0.5, size=26, bold=True, color=CRED_GOLD)
    add_text(slide, prob,    x+0.2, 2.6,  3.7, 0.9, size=13, bold=True)
    add_rect(slide, x+0.2, 3.6, 3.6, 0.03, fill_rgb=CRED_GRAY)
    add_text(slide, insight, x+0.2, 3.75, 3.7, 1.6, size=12, color=CRED_GRAY)
    # Impact tag at card bottom
    add_rect(slide, x+0.2, 5.52, 3.6, 0.26, fill_rgb=RGBColor(0x2A, 0x24, 0x0E))
    add_text(slide, "→ " + impact_tags[i], x+0.3, 5.54, 3.45, 0.22, size=9, color=CRED_GOLD)
add_rect(slide, 0.4, 6.1, 12.5, 1.0, fill_rgb=RGBColor(0x1E, 0x16, 0x06))
add_rect(slide, 0.4, 6.1, 0.06, 1.0, fill_rgb=CRED_GOLD)
add_text(slide, "Delta 4 Theory (Kunal Shah): users switch only when new experience is 4× better. "
         "Bank netbanking = 2/10.  CRED = 6/10.  The switch happens.",
         0.65, 6.2, 12.0, 0.75, size=12, italic=True, color=CRED_GOLD)


# ════════════════════════════════════════════════════════════════
# SECTION DIVIDER — THE PRODUCT
# ════════════════════════════════════════════════════════════════
section_divider("The Product", "Evolution · Features · App UI")


# ════════════════════════════════════════════════════════════════
# SLIDE 7 — PRODUCT EVOLUTION
# ════════════════════════════════════════════════════════════════
slide = new_slide("PRODUCT EVOLUTION", "From bill payment to financial super-app")
add_text(slide, "How the Product Evolved", 0.6, 0.9, 12, 0.7, size=36, bold=True)
timeline = [
    ("2018", "Launch",     "Credit card\nbill payment\n+ CRED Coins"),
    ("2019", "Rewards",    "CRED Store:\ncoins → brand\noffers"),
    ("2020", "Lending",    "CRED Cash\n(credit line)\n+ CRED Stash"),
    ("2021", "Expansion",  "Travel · RentPay\nMint (FD)"),
    ("2022", "Payments",   "CRED Pay (UPI)\n+ Flash (BNPL)"),
    ("2023–24", "Super-app","Garage · Money\n(vehicles, unified\nfinance)"),
]
dot_y = 2.65   # moved up from 3.3 — creates room for content below
add_rect(slide, 0.55, dot_y+0.18, 12.2, 0.05, fill_rgb=CRED_GOLD)
xs = [0.55, 2.6, 4.65, 6.7, 8.75, 10.8]
for i, (yr, lbl, desc) in enumerate(timeline):
    x = xs[i]
    dot = slide.shapes.add_shape(9, Inches(x+0.78), Inches(dot_y+0.1),
                                  Inches(0.3), Inches(0.3))
    dot.fill.solid(); dot.fill.fore_color.rgb = CRED_GOLD; dot.line.fill.background()
    add_text(slide, yr,   x, dot_y-1.0, 2.0, 0.4,  size=12, bold=True, color=CRED_GOLD, align=PP_ALIGN.CENTER)
    add_text(slide, lbl,  x, dot_y-0.55, 2.0, 0.38, size=11, color=CRED_GRAY, align=PP_ALIGN.CENTER)
    add_rect(slide, x+0.15, dot_y+0.55, 1.6, 1.7, fill_rgb=RGBColor(0x22, 0x22, 0x22))
    add_text(slide, desc, x+0.2,  dot_y+0.6,  1.55, 1.6, size=10, align=PP_ALIGN.CENTER)

# Pattern insight strip (fills empty space below timeline boxes)
add_rect(slide, 0.4, 5.1, 12.5, 1.0, fill_rgb=RGBColor(0x22, 0x22, 0x22))
add_rect(slide, 0.4, 5.1, 0.06, 1.0, fill_rgb=CRED_GOLD)
add_text(slide, "THE PATTERN:", 0.62, 5.18, 2.5, 0.35, size=10, bold=True, color=CRED_GOLD)
add_text(slide, "Each feature either (a) deepens engagement with the existing user base, or "
         "(b) monetises their trust in a new category.  CRED is not chasing new users — "
         "it is extracting more value from the same premium cohort.",
         0.62, 5.55, 12.1, 0.48, size=11, color=CRED_WHITE)

# Key milestone highlights bar
milestones = [("2018", "Entry barrier established"), ("2020", "Lending arm launched"), ("2024", "First EBITDA-positive quarter")]
for i, (yr, lbl) in enumerate(milestones):
    x = 0.4 + i * 4.3
    add_rect(slide, x, 6.22, 4.1, 0.7, fill_rgb=RGBColor(0x1E, 0x1A, 0x08))
    add_rect(slide, x, 6.22, 4.1, 0.05, fill_rgb=CRED_GOLD)
    add_text(slide, yr,  x+0.2, 6.3,  1.0, 0.35, size=12, bold=True, color=CRED_GOLD)
    add_text(slide, lbl, x+1.3, 6.3,  2.65, 0.5, size=11)


# ════════════════════════════════════════════════════════════════
# SLIDE 8 — FEATURE ANALYSIS
# ════════════════════════════════════════════════════════════════
slide = new_slide("FEATURE ANALYSIS", "What each feature does and where it struggles")
add_text(slide, "Feature Analysis", 0.6, 0.9, 12, 0.7, size=36, bold=True)
features = [
    ("Bill Payment",  "✅ Strong",    "Solves real pain. Builds monthly habit. Gateway product.",         "Now commoditised — GPay/PhonePe offer the same utility."),
    ("CRED Coins",    "✅⚠️ Mixed",   "Creates Pavlovian loop; drives retention and repeat opens.",       "Coin devaluation erodes trust with power users."),
    ("CRED Store",    "✅ Strong",    "Brands pay premium CPMs to reach high-converting users.",          "Store has become noisy; curation quality has dropped."),
    ("CRED Cash",     "✅ Strong",    "Low-risk lending (750+ CIBIL), high margins, zero extra CAC.",     "RBI BNPL crackdown forced product restructuring 2022."),
    ("CRED Garage",   "✅ Underrated","Vehicle data → insurance upsell → recurring revenue.",            "Low awareness; undermarketed relative to its potential."),
    ("CRED Travel",   "⚠️ Weak",     "Coin redemption play keeps users on platform.",                   "No sustainable edge vs MMT/Ixigo/Goibibo."),
]
col_x = [0.35, 2.55, 3.95, 8.05]
col_w = [2.1, 1.3, 4.0, 4.1]
row_y = 1.85
add_rect(slide, 0.35, row_y, 12.6, 0.38, fill_rgb=CRED_GOLD)
for hdr, cx, cw in zip(["Feature","Signal","What Works","Risk / Gap"], col_x, col_w):
    add_text(slide, hdr, cx, row_y+0.04, cw, 0.35, size=11, bold=True, color=CRED_BLACK)
for i, (feat, sig, works, risk) in enumerate(features):
    ry = row_y + 0.42 + i * 0.74
    bg = RGBColor(0x22,0x22,0x22) if i%2==0 else RGBColor(0x1E,0x1E,0x1E)
    add_rect(slide, 0.35, ry, 12.6, 0.72, fill_rgb=bg)
    add_text(slide, feat,  col_x[0], ry+0.08, col_w[0], 0.6, size=11, bold=True)
    sc = ACCENT_GREEN if "✅" in sig and "⚠️" not in sig else (CRED_GOLD if "⚠️" in sig else ACCENT_RED)
    add_text(slide, sig,   col_x[1], ry+0.08, col_w[1]+0.4, 0.6, size=10, color=sc)
    add_text(slide, works, col_x[2], ry+0.08, col_w[2], 0.6, size=10)
    add_text(slide, risk,  col_x[3], ry+0.08, col_w[3], 0.6, size=10, color=CRED_GRAY)

# Footer summary bar — fills 0.77" gap at bottom
add_rect(slide, 0.35, 6.78, 12.6, 0.58, fill_rgb=RGBColor(0x1E, 0x1E, 0x1E))
summary_pts = [("Core revenue drivers:", "Bill Pay + Lending + Store"), ("Biggest risk:", "Coins devaluation → trust erosion"), ("Hidden gem:", "CRED Garage — insurance upsell flywheel")]
for j, (label, val) in enumerate(summary_pts):
    sx = 0.55 + j * 4.2
    add_rect(slide, sx, 6.81, 0.06, 0.52, fill_rgb=CRED_GOLD)
    add_text(slide, label, sx+0.15, 6.83, 1.6, 0.25, size=9, bold=True, color=CRED_GOLD)
    add_text(slide, val,   sx+0.15, 7.09, 3.9, 0.24, size=9, color=CRED_WHITE)


# ════════════════════════════════════════════════════════════════
# SLIDE 9 — APP UI SHOWCASE
# ════════════════════════════════════════════════════════════════
slide = new_slide("APP UI SHOWCASE", "Key screens — Bills · Lending · Garage")
add_text(slide, "App UI — Key Screens", 0.6, 0.9, 12, 0.7, size=34, bold=True)
# Three phone mockups — reduced width so labels don't overlap
add_img(slide, PHONE_BILLS,  0.8,  1.6, 2.0)   # right edge 2.8
add_img(slide, PHONE_CASH,   5.0,  1.6, 2.0)   # right edge 7.0
add_img(slide, PHONE_GARAGE, 9.15, 1.6, 2.0)   # right edge 11.15

# Labels (phone bottom ≈ 1.6 + 2.0*(580/320) = 1.6 + 3.625 = 5.22)
for x, lbl, desc in [
    (0.75, "My Bills",     "Monthly bill aggregation\nand one-tap payment"),
    (4.95, "CRED Cash",    "Pre-approved instant\ncredit line"),
    (9.1,  "CRED Garage",  "Vehicle management,\nFASTag & insurance"),
]:
    add_text(slide, lbl,  x, 5.38, 3.0, 0.38, size=13, bold=True, color=CRED_GOLD)
    add_text(slide, desc, x, 5.78, 3.0, 0.55, size=11, color=CRED_GRAY)

# PM annotations
add_rect(slide, 0.4, 6.45, 12.5, 0.85, fill_rgb=RGBColor(0x26, 0x26, 0x26))
add_rect(slide, 0.4, 6.45, 0.06, 0.85, fill_rgb=CRED_GOLD)
add_text(slide, "PM note:  Bills screen is clean and high-trust — solves real pain in one tap.  "
         "CRED Cash converts well because the pre-approved limit is shown upfront (no application anxiety).  "
         "Garage is the most undervalued screen — vehicle data compounds into insurance upsell.",
         0.6, 6.52, 12.1, 0.72, size=10, italic=True, color=CRED_GRAY)


# ════════════════════════════════════════════════════════════════
# SECTION DIVIDER — THE BUSINESS
# ════════════════════════════════════════════════════════════════
section_divider("The Business", "Model · Revenue · Market")


# ════════════════════════════════════════════════════════════════
# SLIDE 11 — BUSINESS MODEL
# ════════════════════════════════════════════════════════════════
slide = new_slide("BUSINESS MODEL", "How CRED actually makes money")
add_text(slide, "Business Model", 0.6, 0.9, 12, 0.7, size=36, bold=True)
pillars = [
    ("01", "Advertising\n& Commerce",
     "Brands pay premium CPMs for placement in CRED Store and targeted campaigns. "
     "Higher conversion than generic digital ads due to verified CIBIL-filtered audience.",
     "High margin"),
    ("02", "Lending",
     "CRED Cash, Flash, Stash — instant credit lines and BNPL. "
     "Low NPA risk (750+ CIBIL). Interest income + processing fees. "
     "₹19,000 Cr AUM as of Dec 2024.",
     "Very High margin"),
    ("03", "Transactional\nCommissions",
     "Travel bookings, RentPay, Garage services — CRED earns a commission "
     "on every transaction completed through the platform.",
     "Medium margin"),
]
est_rev = ["~₹500 Cr est. FY24", "₹19,000 Cr AUM · NPA 1.1%", "~₹200 Cr commissions est."]
for i, (num, title, desc, margin) in enumerate(pillars):
    x = 0.45 + i * 4.3
    add_rect(slide, x, 1.85, 4.1, 4.5, fill_rgb=RGBColor(0x22, 0x22, 0x22))
    add_rect(slide, x, 1.85, 4.1, 0.06, fill_rgb=CRED_GOLD)
    add_text(slide, num,   x+0.2, 1.95, 0.8, 0.55, size=28, bold=True, color=CRED_GOLD)
    add_text(slide, title, x+0.2, 2.55, 3.7, 0.7,  size=15, bold=True)
    add_text(slide, desc,  x+0.2, 3.35, 3.7, 2.0,  size=12, color=CRED_GRAY)
    add_rect(slide, x+0.2, 5.42, 3.6, 0.03, fill_rgb=CRED_GOLD)
    add_text(slide, margin,    x+0.2, 5.5,  3.6, 0.28, size=11, bold=True, color=CRED_GOLD)
    add_text(slide, est_rev[i], x+0.2, 5.82, 3.6, 0.28, size=10, color=CRED_GRAY)
add_rect(slide, 0.4, 6.5, 12.5, 0.88, fill_rgb=RGBColor(0x26, 0x26, 0x26))
add_rect(slide, 0.4, 6.5, 0.06, 0.88, fill_rgb=CRED_GOLD)
add_text(slide, "Master play: Bill payment is a FREE acquisition tool → build trust + data → "
         "monetise through lending and commerce.  Contribution margin positive for 9 consecutive quarters by FY24.",
         0.65, 6.6, 12.0, 0.7, size=12, italic=True, color=CRED_GOLD)


# ════════════════════════════════════════════════════════════════
# SLIDE 12 — BY THE NUMBERS
# ════════════════════════════════════════════════════════════════
slide = new_slide("BY THE NUMBERS", "Public financial data FY22–FY24  ·  Sources: Inc42, Entrackr, TechCrunch")
add_text(slide, "By the Numbers", 0.6, 0.9, 12, 0.7, size=36, bold=True)
add_img(slide, REV_CHART, 0.4, 1.7, 6.5, 2.9)

# Insight boxes below the chart (fills empty left-side space)
chart_insights = [
    (CRED_GOLD,  "6.3× revenue growth", "from ₹394 Cr (FY22) to ₹2,473 Cr (FY24)"),
    (ACCENT_RED, "Operating loss narrowing", "₹1,024 Cr (FY23) → ₹609 Cr (FY24) → targeting EBITDA+"),
    (CRED_GOLD,  "Lending is the margin engine", "9 consecutive quarters of positive contribution margin"),
]
for j, (col, heading, detail) in enumerate(chart_insights):
    cy = 4.75 + j * 0.85
    add_rect(slide, 0.4, cy, 6.5, 0.78, fill_rgb=RGBColor(0x22, 0x22, 0x22))
    add_rect(slide, 0.4, cy, 0.06, 0.78, fill_rgb=col)
    add_text(slide, heading, 0.58, cy+0.08, 5.7, 0.28, size=11, bold=True, color=col)
    add_text(slide, detail,  0.58, cy+0.38, 5.7, 0.34, size=10, color=CRED_GRAY)

stats = [
    ("13M",        "Monthly Active Users",    "Flat since Nov 2022 — intentional"),
    ("1.1%",       "NPA on CRED Cash",         "vs 2.7% banking industry avg"),
    ("₹19,000 Cr", "Lending AUM",              "As of Dec 2024"),
    ("9 Qtrs",     "Contribution Margin +ve",  "Path to EBITDA break-even"),
    ("-41%",       "Operating Loss Drop",       "FY23→FY24: ₹1,024Cr → ₹609Cr"),
    ("$942M",      "Total Raised",             "11 rounds · $6.4B peak val."),
]
for i, (num, label, note) in enumerate(stats):
    row, col = i//3, i%3
    x = 7.2 + col * 2.05
    y = 1.7 + row * 2.45
    add_rect(slide, x, y, 1.95, 2.2, fill_rgb=RGBColor(0x22, 0x22, 0x22))
    add_rect(slide, x, y, 1.95, 0.05, fill_rgb=CRED_GOLD)
    add_text(slide, num,   x+0.12, y+0.18, 1.7, 0.7, size=20, bold=True, color=CRED_GOLD)
    add_text(slide, label, x+0.12, y+0.9,  1.7, 0.6, size=10, bold=True)
    add_text(slide, note,  x+0.12, y+1.52, 1.7, 0.62, size=9, color=CRED_GRAY)


# ════════════════════════════════════════════════════════════════
# SLIDE 13 — HONEST ASSESSMENT
# ════════════════════════════════════════════════════════════════
slide = new_slide("HONEST ASSESSMENT", "What's working and what isn't")
add_text(slide, "Honest Assessment", 0.6, 0.9, 12, 0.7, size=36, bold=True)
add_rect(slide, 0.4, 1.85, 6.1, 0.45, fill_rgb=ACCENT_GREEN)
add_text(slide, "  ✓  What's Working", 0.4, 1.85, 6.1, 0.45, size=14, bold=True)
working = ["Trust moat — 750+ CIBIL filter creates self-selecting quality users",
           "Lending business — low NPAs, high margins, zero extra CAC",
           "Design & brand — one of India's best-designed consumer apps",
           "CRED Garage — quietly strong; vehicle data = insurance goldmine",
           "Core retention — bill reminders create reliable monthly re-engagement"]
add_ml(slide, [("✓  "+w, False, CRED_WHITE) for w in working],
       0.4, 2.35, 5.9, 4.5, size=13)
add_rect(slide, 7.0, 1.85, 6.1, 0.45, fill_rgb=ACCENT_RED)
add_text(slide, "  ✗  What's Not Working", 7.0, 1.85, 6.1, 0.45, size=14, bold=True)
risks = ["CRED Coins devaluation — heavy users feel cheated",
         "Feature sprawl — too many tabs, weak discoverability",
         "Travel & Commerce — no edge vs dedicated players",
         "High CAC — brand campaigns expensive to sustain at scale",
         "RBI regulatory risk — lending & BNPL under ongoing scrutiny"]
add_ml(slide, [("✗  "+r, False, CRED_WHITE) for r in risks],
       7.0, 2.35, 5.9, 4.5, size=13)


# ════════════════════════════════════════════════════════════════
# SLIDE 14 — INDIA MARKET OPPORTUNITY
# ════════════════════════════════════════════════════════════════
slide = new_slide("MARKET OPPORTUNITY", "India's premium credit card market")
add_text(slide, "India Market Opportunity", 0.6, 0.9, 12, 0.7, size=34, bold=True)
add_img(slide, INDIA_MKT, 0.4, 1.7, 5.5, 3.5)

opp = [
    ("100M+",    "Credit card holders in India (growing at ~20% YoY)"),
    ("35–40M",   "Premium holders (750+ CIBIL) — CRED's actual TAM"),
    ("13M",      "CRED MAU — ~33% TAM penetration"),
    ("₹8–10T",   "Annual credit card spend in India (FY24)"),
    ("₹2,473 Cr","CRED revenue FY24 on 13M users = ₹190/user/yr"),
    ("10×",      "Revenue/user headroom vs global fintech benchmarks"),
]
add_rect(slide, 6.3, 1.72, 6.6, 0.42, fill_rgb=CRED_GOLD)
add_text(slide, "  Key Market Numbers", 6.3, 1.72, 6.6, 0.42, size=12, bold=True, color=CRED_BLACK)
for i, (num, label) in enumerate(opp):
    ry = 2.18 + i * 0.60           # FIXED: was 0.77 — caused overlap with insight box
    bg = RGBColor(0x22,0x22,0x22) if i%2==0 else RGBColor(0x1E,0x1E,0x1E)
    add_rect(slide, 6.3, ry, 6.6, 0.57, fill_rgb=bg)
    add_text(slide, num,   6.5,  ry+0.08, 1.4, 0.35, size=14, bold=True, color=CRED_GOLD)
    add_text(slide, label, 7.95, ry+0.10, 4.7, 0.38, size=11)

# Insight box now safely below table (last row ends at ~5.78)
add_rect(slide, 0.4, 5.9, 12.5, 0.85, fill_rgb=RGBColor(0x1E, 0x16, 0x06))
add_rect(slide, 0.4, 5.9, 0.06, 0.85, fill_rgb=CRED_GOLD)
add_text(slide, "CRED's opportunity isn't user growth — it's revenue depth per user. "
         "At ₹190/user/yr it is massively undermonetised vs international benchmarks (Amex: ₹8,000+/user/yr). "
         "Lending at scale is the unlock.",
         0.65, 5.98, 12.0, 0.72, size=12, italic=True, color=CRED_GOLD)


# ════════════════════════════════════════════════════════════════
# SECTION DIVIDER — THE COMPETITION
# ════════════════════════════════════════════════════════════════
section_divider("The Competition", "Who's threatening CRED's moat in 2025–26")


# ════════════════════════════════════════════════════════════════
# SLIDE 16 — CLASSIC 4-PLAYER TABLE
# ════════════════════════════════════════════════════════════════
slide = new_slide("COMPETITION", "Mass-market players vs CRED")
add_text(slide, "CRED vs Mass Market", 0.6, 0.9, 12, 0.7, size=34, bold=True)

# Correct logos for THIS slide: PhonePe, Paytm, GPay (the mass-market players being compared)
for lx, logo in [(5.55, LOGO_PHONEPE), (8.2, LOGO_PAYTM), (10.85, LOGO_GPAY)]:
    add_img(slide, logo, lx+0.95, 1.08, 0.55, 0.55)

rows = [
    ("Dimension",       "CRED",              "PhonePe",       "Paytm",           "GPay"),
    ("Primary Utility", "CC Bills + Rewards","UPI Payments",  "Payments + Lend.", "UPI Payments"),
    ("Target User",     "Premium 750+ CIBIL","Mass market",   "Mass market",     "Mass market"),
    ("Rewards",         "Strong (coins)",    "Cashback",      "Cashback",        "Minimal"),
    ("Lending",         "Yes (premium)",     "Yes",           "Yes (aggressive)","No"),
    ("Brand",           "Premium/exclusive", "Utility",       "Declining",       "Google-backed"),
    ("Design",          "Best-in-class",     "Good",          "Average",         "Clean"),
]
cx = [0.35, 3.0, 5.65, 8.3, 10.95]
cw = [2.6, 2.6, 2.6, 2.6, 2.3]
for i, row in enumerate(rows):
    ry = 1.85 + i * 0.56
    if i == 0:
        add_rect(slide, 0.35, ry, 12.6, 0.56, fill_rgb=CRED_GOLD)
        tc, bld = CRED_BLACK, True
    else:
        bg = RGBColor(0x22,0x22,0x22) if i%2==1 else RGBColor(0x1E,0x1E,0x1E)
        add_rect(slide, 0.35, ry, 12.6, 0.56, fill_rgb=bg)
        tc, bld = CRED_WHITE, False
    for j, (cell, c, w) in enumerate(zip(row, cx, cw)):
        if j == 1 and i > 0:
            add_rect(slide, c-0.05, ry, w+0.1, 0.56, fill_rgb=RGBColor(0x2A,0x26,0x1A))
            cc = CRED_GOLD
        else:
            cc = tc if i > 0 else CRED_BLACK
        add_text(slide, cell, c, ry+0.1, w, 0.42, size=11, bold=(bld or j==0), color=cc)
# Market share visual bar (UPI 2025 data)
add_rect(slide, 0.35, 6.43, 12.6, 0.58, fill_rgb=RGBColor(0x1A, 0x1A, 0x1A))
add_text(slide, "UPI Market Share 2025:", 0.5, 6.5, 2.4, 0.28, size=9, bold=True, color=CRED_GOLD)
mshare = [("PhonePe 47%", 47, "4A90D9"), ("GPay 37%", 37, "5DA0E8"),
          ("Paytm 10%", 10, "888888"), ("CRED Pay ~3%", 3, "C9A84C")]
xb = 3.1
for label, pct, col in mshare:
    bw = pct * 0.115
    bc = RGBColor(*[int(col[i:i+2], 16) for i in (0,2,4)])
    add_rect(slide, xb, 6.49, bw, 0.24, fill_rgb=bc)
    add_text(slide, label, xb, 6.76, bw+0.5, 0.2, size=7.5, color=CRED_GRAY)
    xb += bw + 0.22

add_text(slide, "CRED's defensible position: only platform that has earned trust of India's creditworthy, high-income segment at scale.",
         0.5, 7.12, 12.3, 0.3, size=10, italic=True, color=CRED_GRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 17 — COMPETITIVE THREAT 2025-26 (with logos)
# ════════════════════════════════════════════════════════════════
slide = new_slide("COMPETITIVE THREAT LANDSCAPE 2025–26", "The real challengers to CRED's moat")
add_text(slide, "Fintech Threat Landscape", 0.6, 0.9, 12, 0.7, size=34, bold=True)

comp_data = [
    (LOGO_OC,     "OneCard",   "Urban prof.\n21–35",
     "Own issuer\n(SBM/BoB)",   "5× pts on\ntop 2 categories",
     "HIGH", "Same demo, owns the card relationship"),
    (LOGO_SLICE,  "Slice SFB", "18–30, credit\nnewcomers",
     "Full bank post\nNESFB merger",  "Cashback + savings\naccount bundle",
     "MEDIUM", "Captures users before CRED-eligibility"),
    (LOGO_UNI,    "Uni Cards", "Premium prof,\nhigh spenders",
     "Co-branded\npartner banks","Pay 1/3 monthly\n(split payments)",
     "LOW", "₹95 Cr revenue — struggling to scale"),
    (LOGO_JUP,    "Jupiter",   "Digital-first\nmillennials",
     "Neobank,\nFederal Bank","Smart pots +\nRuPay UPI card",
     "HIGH", "Full financial OS for the same user"),
    (LOGO_FI,     "Fi Money",  "Tech professionals\nhigh income",
     "Neobank,\nFederal Bank","Smart deposits +\npremium credit line",
     "HIGH", "Lifestyle + credit competitor"),
    (LOGO_SCAPIA, "Scapia",    "Travel lovers,\nfrequent flyers",
     "Co-branded\nFederal Bank","20% coins travel,\nzero forex, zero fee",
     "MED-HIGH", "Steals CRED Travel segment users"),
]

col_xs = [0.35, 1.55, 3.1, 5.1, 7.1, 9.0, 11.0]
col_ws = [1.05, 1.5, 1.9, 1.9, 1.8, 1.9, 2.2]
hdrs   = ["", "Company", "Target User", "Card Model", "Rewards", "Threat", "Why"]

add_rect(slide, 0.35, 1.78, 12.6, 0.38, fill_rgb=CRED_GOLD)
for hdr, cx2, cw2 in zip(hdrs, col_xs, col_ws):
    add_text(slide, hdr, cx2, 1.82, cw2, 0.32, size=10, bold=True, color=CRED_BLACK)

for i, (logo_path, name, user, card, reward, threat, why) in enumerate(comp_data):
    ry = 2.2 + i * 0.84
    bg = RGBColor(0x22,0x22,0x22) if i%2==0 else RGBColor(0x1E,0x1E,0x1E)
    add_rect(slide, 0.35, ry, 12.6, 0.82, fill_rgb=bg)
    add_img(slide, logo_path, col_xs[0]+0.05, ry+0.1, 0.62, 0.62)
    threat_c = ACCENT_RED if "HIGH" == threat else (ACCENT_AMBER if "MED" in threat else ACCENT_GREEN)
    for j, (val, cx2, cw2) in enumerate(zip([name, user, card, reward, threat, why],
                                             col_xs[1:], col_ws[1:])):
        c = threat_c if j == 4 else CRED_WHITE
        add_text(slide, val, cx2, ry+0.1, cw2, 0.68, size=10,
                 bold=(j==0 or j==4), color=c)


# ════════════════════════════════════════════════════════════════
# SECTION DIVIDER — PRODUCT THINKING
# ════════════════════════════════════════════════════════════════
section_divider("Product Thinking", "Opportunities · Pain Points · Onboarding · North Star")


# ════════════════════════════════════════════════════════════════
# SLIDE 19 — WHAT I WOULD BUILD NEXT
# ════════════════════════════════════════════════════════════════
slide = new_slide("PRODUCT OPPORTUNITIES", "What I would build next")
add_text(slide, "What I Would Build Next", 0.6, 0.9, 12, 0.7, size=34, bold=True)
opps = [
    ("01", "CRED Score Dashboard",
     "Problem: Users open CRED once/month and leave. "
     "No daily reason to return.\n\n"
     "Build: Proprietary health index — bill behaviour, credit utilisation, spend patterns, debt-to-income.\n\n"
     "Why: Daily open reason, new data asset, natural CRED Cash upsell."),
    ("02", "CRED for Business",
     "Problem: Self-employed & founder users (large CRED cohort) "
     "can't separate personal vs business card spend.\n\n"
     "Build: 'Business Mode' — tag transactions, generate spend reports, export to accounting.\n\n"
     "Why: Low build cost, high perceived value, opens B2B lane."),
    ("03", "Proactive Spend Intelligence",
     "Problem: Users know their balance, not their behaviour.\n\n"
     "Build: Monthly AI 'Spend Story' — anomalies + one recommended action "
     "(e.g. 'Switch this sub to Card B for 2× rewards').\n\n"
     "Why: Increases DAU, deepens trust, upsells CRED Cash or Mint."),
]
rice = ["RICE ~520 · HIGH", "RICE ~390 · MED", "RICE ~280 · MED"]
for i, (num, title, desc) in enumerate(opps):
    x = 0.4 + i * 4.3
    add_rect(slide, x, 1.85, 4.1, 5.0, fill_rgb=RGBColor(0x22, 0x22, 0x22))
    add_rect(slide, x, 1.85, 4.1, 0.06, fill_rgb=CRED_GOLD)
    add_text(slide, num,   x+0.2, 1.95, 0.8, 0.5,  size=24, bold=True, color=CRED_GOLD)
    add_text(slide, title, x+0.2, 2.5,  3.7, 0.55,  size=14, bold=True)
    add_text(slide, desc,  x+0.2, 3.1,  3.7, 3.3,  size=11, color=CRED_GRAY)
    # RICE badge at card bottom
    add_rect(slide, x+0.2, 6.4, 3.6, 0.28, fill_rgb=RGBColor(0x2A, 0x24, 0x0E))
    add_text(slide, rice[i], x+0.3, 6.42, 3.45, 0.24, size=10, bold=True, color=CRED_GOLD)

# Priority order strip below cards
add_rect(slide, 0.4, 6.82, 12.5, 0.48, fill_rgb=RGBColor(0x22, 0x22, 0x22))
add_text(slide, "Recommended sequence:  CRED Score Dashboard  →  Spend Intelligence  →  Business Mode",
         0.58, 6.88, 12.2, 0.36, size=11, color=CRED_GOLD)


# ════════════════════════════════════════════════════════════════
# SLIDE 20 — PRODUCT OPPORTUNITY BRIEF
# ════════════════════════════════════════════════════════════════
slide = new_slide("PRODUCT OPPORTUNITY BRIEF", "CRED Score Dashboard — Full Brief")
add_text(slide, "Product Brief: CRED Score Dashboard", 0.6, 0.9, 12, 0.65, size=28, bold=True)

def brief_box(slide, title, content, x, y, w, h, tc=CRED_BLACK, bc=CRED_GOLD):
    add_rect(slide, x, y, w, 0.38, fill_rgb=bc)
    add_text(slide, "  "+title, x, y, w, 0.38, size=11, bold=True, color=tc)
    add_rect(slide, x, y+0.38, w, h-0.38, fill_rgb=RGBColor(0x22, 0x22, 0x22))
    return y+0.42

y0 = brief_box(slide, "Problem Statement", "", 0.4, 1.68, 6.1, 1.25)
add_text(slide, "Users open CRED once/month to pay bills then leave. CIBIL is a lagging indicator — "
         "tells users what happened, not what to do. No competitor offers a proactive financial health OS for premium users.",
         0.55, y0, 5.85, 0.82, size=11, color=CRED_GRAY)

y1 = brief_box(slide, "Target User", "", 0.4, 3.05, 6.1, 0.85)
add_text(slide, "Existing CRED members who pay bills regularly but haven't activated CRED Cash — "
         "the pre-lending conversion cohort.",
         0.55, y1, 5.85, 0.52, size=11, color=CRED_GRAY)

y2 = brief_box(slide, "MVP Features (MoSCoW)", "", 0.4, 4.0, 6.1, 2.85)
mvp = [("M  CRED Health Score (proprietary index)",   True,  CRED_WHITE),
       ("M  Credit utilisation trend (30/60/90 days)", True,  CRED_WHITE),
       ("S  Spend category breakdown + benchmarks",    False, CRED_WHITE),
       ("S  Bill payment streak & on-time history",    False, CRED_WHITE),
       ("C  'Improve your score' action cards",        False, CRED_GRAY),
       ("W  Competitor card comparison engine",        False, CRED_GRAY)]
add_ml(slide, mvp, 0.55, y2, 5.85, 2.35, size=11)

# Right column
y3 = brief_box(slide, "Success Metrics", "", 6.7, 1.68, 6.1, 1.7)
add_ml(slide, [
    ("PRIMARY", True, CRED_GOLD),
    ("• DAU/MAU ratio  (target: +8pp in 90 days)", False, CRED_WHITE),
    ("• CRED Cash activation rate among dashboard users", False, CRED_WHITE),
    ("GUARDRAIL", True, CRED_GOLD),
    ("• Session length ≤ 10% increase (avoid fatigue)", False, CRED_WHITE),
    ("• CIBIL pull consent rate ≥ 80%", False, CRED_WHITE),
], 6.85, y3, 5.85, 1.28, size=11)

y4 = brief_box(slide, "RICE Score", "", 6.7, 3.5, 6.1, 1.28)
add_ml(slide, [
    ("Reach: 13M MAU (100% eligible)",        False, CRED_WHITE),
    ("Impact: HIGH — lending funnel entry",    False, CRED_WHITE),
    ("Confidence: 80%  ·  Effort: M (8–12 wks)", False, CRED_WHITE),
    ("→  RICE ~520  ·  HIGH PRIORITY",         True,  CRED_GOLD),
], 6.85, y4, 5.85, 1.0, size=11)

y5 = brief_box(slide, "Top 3 Risks", "", 6.7, 4.9, 6.1, 1.95, tc=CRED_WHITE, bc=ACCENT_RED)
add_ml(slide, [
    "1.  RBI data localisation rules may restrict real-time CIBIL pulls",
    "2.  Users may distrust a proprietary score vs official CIBIL",
    "3.  Score gamification may incentivise wrong financial behaviours",
], 6.85, y5, 5.85, 1.5, size=11)


# ════════════════════════════════════════════════════════════════
# SLIDE 21 — USER PAIN POINTS
# ════════════════════════════════════════════════════════════════
slide = new_slide("USER PAIN POINTS", "Synthesised from 1★ & 2★ app reviews — Play Store, Trustpilot, MouthShut")
add_text(slide, "Top 5 User Pain Points", 0.6, 0.9, 12, 0.65, size=34, bold=True)
pain_points = [
    ("01", "CRED Coins Devaluation",             "HIGH",
     "Coins worth less each cycle. Power users with large balances feel cheated.",
     "\"50,000 coins and nothing redeemable for more than ₹50\"",
     "Transparent tier system with locked-in redemption rates"),
    ("02", "No Phone Support / Slow Resolution", "HIGH",
     "No care number. Bot fails on complex issues. Email takes 5–10 days.",
     "\"Multiple emails sent, no response. Money stuck for 3 weeks\"",
     "Live chat for payment issues; 4-hr SLA on transaction failures"),
    ("03", "Payment Failures & Refund Delays",   "HIGH",
     "Transactions marked 'successful' but not credited. Refunds take 7–21 days.",
     "\"Refund initiated after 45 days of follow-up — unacceptable\"",
     "Real-time bank confirmation + auto-refund SLA of 3 days"),
    ("04", "Notification Spam & App Clutter",    "MEDIUM",
     "5+ daily promo notifications. Bill payment (core) is buried in tabs.",
     "\"I get 5 notifications a day about Store deals I never asked for\"",
     "Opt-in notifications by category; surface bill CTA prominently"),
    ("05", "Store Product Quality & MRP Mismatch","MEDIUM",
     "Products delivered with MRP higher than advertised price. Mystery box issues.",
     "\"MRP on product was higher — doesn't feel like a discount at all\"",
     "Verified MRP display; buyer protection policy for Store orders"),
]
for i, (num, title, freq, body, quote, fix) in enumerate(pain_points):
    y = 1.9 + i * 1.06
    add_rect(slide, 0.35, y, 12.6, 1.02, fill_rgb=RGBColor(0x22, 0x22, 0x22))
    fc = ACCENT_RED if freq == "HIGH" else CRED_GOLD
    add_rect(slide, 0.35, y, 0.07, 1.02, fill_rgb=fc)
    add_text(slide, num,   0.55, y+0.05, 0.55, 0.4, size=14, bold=True, color=fc)
    add_text(slide, title, 1.15, y+0.06, 2.9,  0.38, size=12, bold=True)
    add_text(slide, freq,  4.1,  y+0.09, 0.95, 0.28, size=10, bold=True, color=fc)
    add_text(slide, body,  1.15, y+0.52, 5.0,  0.48, size=10, color=CRED_GRAY)
    add_rect(slide, 6.4, y+0.05, 3.2, 0.92, fill_rgb=RGBColor(0x1A, 0x1A, 0x1A))
    add_text(slide, quote, 6.52, y+0.1,  3.0, 0.82, size=9, italic=True, color=CRED_GOLD)
    add_text(slide, "Fix: "+fix, 9.7, y+0.1, 3.05, 0.88, size=9, color=CRED_WHITE)


# ════════════════════════════════════════════════════════════════
# SLIDE 22 — ONBOARDING FLOW
# ════════════════════════════════════════════════════════════════
slide = new_slide("ONBOARDING FLOW", "From app download to first bill payment")
add_text(slide, "Onboarding Flow Analysis", 0.6, 0.9, 12, 0.65, size=32, bold=True)

# Context banner — fills gap between title and flow start
add_rect(slide, 0.4, 1.68, 12.5, 0.48, fill_rgb=RGBColor(0x22, 0x22, 0x22))
add_rect(slide, 0.4, 1.68, 0.06, 0.48, fill_rgb=CRED_GOLD)
add_text(slide, "ENTRY CONTEXT:", 0.62, 1.74, 2.0, 0.28, size=9, bold=True, color=CRED_GOLD)
add_text(slide, "~35–40M Indians qualify (750+ CIBIL).  "
         "CRED's onboarding must preserve exclusivity while minimising drop-off in this pre-qualified cohort.",
         2.72, 1.74, 9.9, 0.32, size=9, color=CRED_GRAY)

steps = [
    ("App\nDiscovery",  "Invite from friend\nor IPL ad",
     "✅ Exclusivity creates\naspirations pre-install",
     "⚠️ Invite gate\nblocks organic growth"),
    ("Download\n& Open",  "Play Store install;\nCRED brand splash",
     "✅ Premium first\nimpression on open",
     "⚠️ No value preview\nbefore sign-up"),
    ("OTP\nVerify",  "Mobile → OTP;\n< 30 sec",
     "✅ Frictionless;\nwell-designed",
     "⚠️ None — step\nis clean"),
    ("CIBIL\nCheck",  "PAN + DOB;\nscore auto-fetched",
     "✅ Instant; user sees\nscore for free",
     "🔴 HIGH DROP-OFF\nFear of hard inquiry"),
    ("Card\nLinking",  "Manual entry or\nauto-detect",
     "✅ Supports all\nmajor networks",
     "⚠️ Manual entry\ntedious; auto fails"),
    ("First\nPayment",  "Dashboard shows due;\none-tap pay",
     "✅ Clean UI;\ninstant coins earned",
     "⚠️ Coin value opaque\nbefore first pay"),
]
dot_y = 3.55   # shifted down slightly to fit context banner above
add_rect(slide, 0.5, dot_y+0.2, 12.2, 0.04, fill_rgb=CRED_GOLD)
sxs = [0.4, 2.45, 4.5, 6.55, 8.6, 10.65]
for i, (step, detail, good, bad) in enumerate(steps):
    x = sxs[i]
    is_critical = "🔴" in bad
    dot_c = ACCENT_RED if is_critical else CRED_GOLD
    dot = slide.shapes.add_shape(9, Inches(x+0.73), Inches(dot_y+0.1),
                                  Inches(0.3), Inches(0.3))
    dot.fill.solid(); dot.fill.fore_color.rgb = dot_c; dot.line.fill.background()
    add_rect(slide, x+0.25, dot_y-1.3, 1.6, 1.22, fill_rgb=RGBColor(0x22, 0x22, 0x22))
    if is_critical:
        add_rect(slide, x+0.25, dot_y-1.3, 1.6, 0.05, fill_rgb=ACCENT_RED)
    add_text(slide, step,   x+0.3, dot_y-1.28, 1.5, 0.5, size=10, bold=True,
             color=CRED_GOLD if not is_critical else ACCENT_RED, align=PP_ALIGN.CENTER)
    add_text(slide, detail, x+0.3, dot_y-0.75, 1.5, 0.5, size=9,
             color=CRED_GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, good, x+0.1, dot_y+0.65, 1.9, 0.78, size=9,
             color=ACCENT_GREEN, align=PP_ALIGN.CENTER)
    add_text(slide, bad,  x+0.1, dot_y+1.5,  1.9, 0.82, size=9,
             color=ACCENT_RED if is_critical else CRED_GOLD, align=PP_ALIGN.CENTER)
# Drop-off risk callout — fills gap between flow and PM fix
add_rect(slide, 0.4, 5.78, 12.5, 0.85, fill_rgb=RGBColor(0x1A, 0x0E, 0x0E))
add_rect(slide, 0.4, 5.78, 0.07, 0.85, fill_rgb=ACCENT_RED)
add_text(slide, "CRITICAL DROP-OFF:", 0.6, 5.85, 2.6, 0.3, size=10, bold=True, color=ACCENT_RED)
add_text(slide, "CIBIL check step — fear of hard inquiry causes an estimated 15–25% of qualified users to abandon. "
         "Fix: display 'soft check only — your score won't change' badge prominently before PAN entry. "
         "Reduce form from 4 fields to 2 (PAN + consent).",
         3.35, 5.82, 9.35, 0.78, size=10, color=CRED_WHITE)

add_rect(slide, 0.4, 6.75, 12.5, 0.58, fill_rgb=RGBColor(0x26, 0x26, 0x26))
add_rect(slide, 0.4, 6.75, 0.06, 0.58, fill_rgb=CRED_GOLD)
add_text(slide, "PM Fixes:  (1) 'Soft check' badge at CIBIL step  ·  "
         "(2) CAMS/CVL one-click card import  ·  "
         "(3) Show coin value before first payment to anchor expectations",
         0.65, 6.82, 12.1, 0.44, size=10, italic=True, color=CRED_GOLD)


# ════════════════════════════════════════════════════════════════
# SLIDE 23 — NORTH STAR METRIC
# ════════════════════════════════════════════════════════════════
slide = new_slide("NORTH STAR METRIC", "What CRED should optimise for")
add_text(slide, "North Star Metric", 0.6, 0.9, 12, 0.65, size=34, bold=True)
candidates = [
    ("Candidate 1",
     "Monthly Bill\nPayment GMV",
     "Measures core utility. Easy to track. Grows with card adoption in India.",
     "Doesn't capture monetisation. GMV can grow while business stagnates. "
     "No signal on lending or commerce revenue.",
     False),
    ("Candidate 2",
     "Revenue Per\nActive User (RPAU)",
     "Captures monetisation efficiency. Aligns with lending + commerce strategy. "
     "Forces product to deepen value for existing fixed user base.",
     "Could incentivise short-term squeezing. "
     "Needs guardrails on NPA and churn rate.",
     True),
    ("Candidate 3",
     "Lending AUM Per\nEligible User",
     "Directly tracks the highest-margin business. "
     "Metric CRED's investors care most about.",
     "Too narrow — ignores 80%+ of users not using CRED Cash. "
     "Neglects the trust layer that makes lending work.",
     False),
]
for i, (label, title, pros, cons, is_rec) in enumerate(candidates):
    x = 0.4 + i * 4.3
    border_c = CRED_GOLD if is_rec else CRED_GRAY
    if is_rec:
        add_rect(slide, x, 1.85, 4.1, 5.1,
                 fill_rgb=RGBColor(0x25,0x22,0x12),
                 line_rgb=CRED_GOLD, line_pt=1.5)
    else:
        add_rect(slide, x, 1.85, 4.1, 5.1, fill_rgb=RGBColor(0x22,0x22,0x22))
    add_rect(slide, x, 1.85, 4.1, 0.06, fill_rgb=border_c)
    add_text(slide, label, x+0.2, 1.93, 3.7, 0.3, size=10, color=CRED_GRAY)
    add_text(slide, title, x+0.2, 2.27, 3.7, 0.72, size=15, bold=True,
             color=CRED_GOLD if is_rec else CRED_WHITE)
    add_rect(slide, x+0.2, 3.1, 3.5, 0.03, fill_rgb=ACCENT_GREEN)
    add_text(slide, "Why it works:", x+0.2, 3.18, 3.7, 0.3, size=10, bold=True, color=ACCENT_GREEN)
    add_text(slide, pros, x+0.2, 3.52, 3.7, 1.28, size=10)
    add_rect(slide, x+0.2, 4.9, 3.5, 0.03, fill_rgb=ACCENT_RED)
    add_text(slide, "Trade-off:", x+0.2, 4.98, 3.7, 0.3, size=10, bold=True, color=ACCENT_RED)
    add_text(slide, cons, x+0.2, 5.32, 3.7, 0.9, size=10, color=CRED_GRAY)
    verdict = ("✓  RECOMMENDED" if is_rec else "✗  Reject")
    add_text(slide, verdict, x+0.2, 6.5, 3.7, 0.3, size=10, bold=True,
             color=CRED_GOLD if is_rec else CRED_GRAY)
add_rect(slide, 0.4, 7.0, 12.5, 0.42, fill_rgb=RGBColor(0x26, 0x26, 0x26))
add_text(slide, "Supporting metrics laddering up to RPAU:  "
         "① Monthly Bill GMV  (core health)   "
         "② Lending AUM Growth  (monetisation depth)   "
         "③ CRED Store GMV per Active User  (commerce engagement)",
         0.6, 7.05, 12.2, 0.34, size=11, color=CRED_GOLD)


# ════════════════════════════════════════════════════════════════
# SLIDE 24 — PM TAKEAWAYS
# ════════════════════════════════════════════════════════════════
slide = new_slide("PM TAKEAWAYS", "Lessons for product builders")
add_text(slide, "PM Takeaways", 0.6, 0.9, 12, 0.7, size=36, bold=True)
takeaways = [
    ("Constraint as strategy",
     "The 750+ CIBIL rule is a product decision, not just marketing. "
     "It sets the quality of every downstream feature — NPA rates, brand perception, ad CPMs. Constraints can be moats."),
    ("Nail the core loop first",
     "Bill payment → coins → store is a tight, validated loop. "
     "CRED earned the right to expand into lending and travel only after this worked. Premature expansion dilutes trust."),
    ("Monetisation follows trust",
     "CRED didn't launch lending on day one. It waited until users trusted the brand "
     "enough to borrow from it. The sequencing of trust-building before monetisation is deliberate strategy."),
    ("Design is a product feature",
     "In a world of functional parity (everyone does bill payments), "
     "CRED's UI and brand voice are genuine retention drivers — not a vanity metric."),
]
for i, (title, body) in enumerate(takeaways):
    row, col = i//2, i%2
    x = 0.4 + col*6.5
    y = 1.85 + row*2.55
    add_rect(slide, x, y, 6.2, 2.38, fill_rgb=RGBColor(0x22, 0x22, 0x22))
    add_rect(slide, x, y, 0.07, 2.38, fill_rgb=CRED_GOLD)
    add_text(slide, title, x+0.28, y+0.18, 5.75, 0.45, size=14, bold=True, color=CRED_GOLD)
    add_text(slide, body,  x+0.28, y+0.72, 5.75, 1.55, size=12, color=CRED_GRAY)


# ════════════════════════════════════════════════════════════════
# SLIDE 25 — END CARD
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=RGBColor(0x0E, 0x0E, 0x0E))
add_rect(slide, 0, 0, 13.33, 0.08, fill_rgb=CRED_GOLD)
add_rect(slide, 0, 7.42, 13.33, 0.08, fill_rgb=CRED_GOLD)
add_rect(slide, 0, 0, 0.12, 7.5, fill_rgb=CRED_GOLD)
add_img(slide, CRED_LOGO, 9.2, 2.5, 3.6, 1.35)
add_text(slide, "Shraddha Singh", 0.8, 2.3, 9, 1.1, size=48, bold=True)
add_text(slide, "Product Manager  ·  Fintech & Enterprise SaaS",
         0.8, 3.52, 9, 0.65, size=18, color=CRED_GOLD)
add_rect(slide, 0.8, 4.38, 4.5, 0.04, fill_rgb=CRED_GOLD)
add_text(slide, "shraddhasingh8968@gmail.com",
         0.8, 4.55, 6, 0.38, size=13, color=CRED_GRAY)
add_text(slide, "linkedin.com/in/shraddha-singh-9149a4172",
         0.8, 4.98, 7, 0.38, size=13, color=CRED_GRAY)
add_rect(slide, 0.8, 5.65, 8, 0.04, fill_rgb=RGBColor(0x40, 0x40, 0x40))
add_text(slide, "PM Portfolio  ·  CRED Product Teardown  ·  Feb 2026",
         0.8, 5.8, 8, 0.38, size=11, color=RGBColor(0x60, 0x60, 0x60))


# ── Save ────────────────────────────────────────────────────────
OUT = "/Users/utkarshsingh/Downloads/workspace/second-brain/portfolio/teardowns/CRED-Teardown-Shraddha-Singh.pptx"
prs.save(OUT)
print(f"\n✅  Saved → {OUT}")
print(f"   Slides: {len(prs.slides)}")
